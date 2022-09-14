# -*-coding:utf-8-*-
# @Time    : 2020/10/17 17:41
# @Author  : Lowe.li
# @Email   : Lowe.li@aishu.cn
import sys
import os
import base64
import gc
import traceback
import multiprocessing
from Crypto.PublicKey import RSA
from Crypto.Cipher import PKCS1_v1_5 as Cipher_pkcsl_v1_5
from concurrent.futures import ThreadPoolExecutor, wait, as_completed
from requests.adapters import HTTPAdapter
import json
import requests
import queue
import json
import logging
from collections import defaultdict
from configparser import ConfigParser
import threading
import re
import time
import datetime
import io
import re
import docx
import pptx
import xlrd
import pdfplumber
import pandas as pd
from itertools import product
from dao.dsm_dao import dsm_dao
import re
from dao.otl_dao import otl_dao
from third_party_service.anyshare.token import asToken
from utils.CommonUtil import commonutil
from c_extend.tools import split_str
from utils.ConnectUtil import mongoClient as MongoClientAD

requests.packages.urllib3.disable_warnings()
sys.path.append(os.path.abspath("../"))


class AccountManager(object):
    def __init__(self,config_file):
        self.config_file=config_file
        self.Config = ConfigParser()
        self.Config.read(self.config_file, "utf-8")

    # rsa加密
    def RSAEndcode(self, msg):
        msg = bytes(msg, encoding="utf8")
        rsakey = RSA.importKey(self.Config['pub_key'].get('pub_key'))
        cipher = Cipher_pkcsl_v1_5.new(rsakey)
        cipher_decode = base64.b64encode(cipher.encrypt(msg)).decode()
        return cipher_decode

    def login(self):
        account=self.Config['loginconfig'].get('account')
        password=self.Config['loginconfig'].get('password')
        params = {'account': account, 'password': self.RSAEndcode(password)}
        r = requests.post(self.Config['url_path'].get('Auth'), data=json.dumps(params), verify=False)
        rs_json = r.json()
        return rs_json

class Node():
    def __init__(self, item):
        self.docid = item["docid"]
        self.name = item["name"]
        self.creator = item["creator"]
        self.editor = item["editor"]
        self.create_time = item["create_time"]
        self.path = item["path"]
        self.rev = item["rev"]
        self.child_dir = []
        self.child_file = []
        self.parent = None

    def add_child(self, child):
        if not isinstance(child, Node):
            self.child_file.append(child)
        else:
            self.child_dir.append(child)
            child.parent = self
def get_api_response(payload, as7json, url):
    ds_auth = as7json['ds_auth']
    print('开始获取AS token', __file__, 'get_api_response')
    ret_token, token_id = asToken.get_token(ds_auth)
    headers = {
        'Content-Type': 'application/json',
        'Authorization': 'Bearer {}'.format(token_id)
    }
    response = requests.request("POST", url, headers=headers, data=json.dumps(payload), verify=False)
    return json.loads(response.text)

class Tree():
    def __init__(self):
        pass #base class

    def add(self, node):
        pass #base class

    def dftraverse(self, node):
        pass #base class


class AnyshareFileTree(Tree):
    def __init__(self, item, as7_json, url, path_url):
        super(AnyshareFileTree, self).__init__()
        self.item = item
        self.url = url
        self.as7_json = as7_json
        self.path_url = path_url
        self.root_node = Node(item)
        self.file_info_dict = {}

    # 深度优先创建树
    def bulid(self):
        def _iter_tree(cur_node, item):
            item_path = get_api_response({"docid": item["docid"]}, self.as7_json, self.path_url).get("namepath")
            if "root" in item:
                self.file_info_dict[item["docid"]] = \
                    (item["docid"], item["name"], item["creator"], item["editor"], item["create_time"],item["rev"], item_path,item["root"], 1)
            else:
                self.file_info_dict[item["docid"]] = \
                    (item["docid"], item["name"], item["creator"], item["editor"], item["create_time"], item["rev"],item_path, 1)
            payload = {"docid": item["docid"], "sort": "asc", "by": "name"}
            child = get_api_response(payload, self.as7_json, self.url)
            child_dir_list = child.get("dirs")
            child_file_list = child.get("files")
            if child_dir_list:
                for child_dir in child_dir_list:
                    child_dir["path"] = get_api_response({"docid": child_dir["docid"]}, self.as7_json, self.path_url).get("namepath")
                    child_dir_node = Node(child_dir)
                    cur_node.add_child(child_dir_node)
                    _iter_tree(child_dir_node, child_dir)
            if child_file_list:
                for child_file in child_file_list:
                    child_path = get_api_response({"docid": child_file["docid"]}, self.as7_json, self.path_url).get("namepath")
                    self.file_info_dict[child_file["docid"]] = \
                        (child_file["docid"], child_file["name"], child_file["creator"], child_file["editor"],
                         child_file["create_time"],child_file["rev"], child_path, 0)
                    cur_node.add_child(child_file["name"])

        _iter_tree(self.root_node, self.item)

    # 深度优先遍历树
    def dftraverse(self, node=None, func=None):
        if not node:
            node = self.root_node
        if not isinstance(node, Node):
            return
        if func:
            func(node)

        for n in node.child_dir:
            self.dftraverse(n)
class GetAnyshareFileTree():
    def __init__(self, as7_json, entrance_dict,config_path):

        self.login_json = None
        self.as7_json = as7_json
        self.entrance_dict = entrance_dict
        self.config_path = config_path
        self.fileInfo_file = "fileInfo.txt"
        self.Config = self.get_config()
        ip = self.as7_json["ds_address"]
        port = self.as7_json["ds_port"]
        if "http" in ip:
            self.ip_and_port = ip + ":" + str(port)
        else:
            self.ip_and_port = "https://" + ip + ":" + str(port)

    def get_config(self):
        config = ConfigParser()
        config.read(self.config_path, "utf-8")
        return config


    def login(self):
        self.login_json = AccountManager(self.config_path).login()

    def save_file_info(self, file_info_dict):
        with open(self.fileInfo_file, "a", encoding="utf-8") as f:
            for docid in file_info_dict:
                f.write("##".join(map(lambda x: str(x), file_info_dict[docid])))
                f.write("\n")

    def get_file_tree(self):
        # 文件树入口
        docid = self.entrance_dict["docid"]
        name = self.entrance_dict["name"]
        creator = self.entrance_dict["creator"]
        create_time = self.entrance_dict["create_time"]
        editor = self.entrance_dict["editor"]
        rev = self.entrance_dict["rev"]
        path = get_api_response({"docid": docid}, self.as7_json,self.ip_and_port + self.Config['url_path'].get('convertPath'))
        new_item_dict = {"docid": docid, "name": name, "creator": creator,"editor": editor, "create_time": create_time, "path": path , "rev":rev,"root":1}
        url = self.ip_and_port + self.Config['url_path'].get('fileListPath')
        path_url = self.ip_and_port + self.Config['url_path'].get('convertPath')
        AFT = AnyshareFileTree(new_item_dict, self.as7_json, url, path_url)
        AFT.bulid()
        return AFT.file_info_dict
def get_gnsInfo_dict(gnsInfo_dict):
    gnsInfo_new_dict = defaultdict(dict)
    for gns in gnsInfo_dict:
        if len(gnsInfo_dict[gns]) == 8:
            _, name, creator, editor, create_time,rev, path, gns_type = gnsInfo_dict[gns]
            gnsInfo_new_dict[gns]["gns"] = gns
            gnsInfo_new_dict[gns]["name"] = name
            gnsInfo_new_dict[gns]["creator"] = creator
            gnsInfo_new_dict[gns]["editor"] = editor
            gnsInfo_new_dict[gns]["create_time"] = create_time
            gnsInfo_new_dict[gns]["path"] = path
            gnsInfo_new_dict[gns]["type"] = gns_type
            gnsInfo_new_dict[gns]["rev"] = rev
        elif len(gnsInfo_dict[gns]) == 9:
            _, name, creator, editor, create_time, rev, path,root ,gns_type  = gnsInfo_dict[gns]
            gnsInfo_new_dict[gns]["gns"] = gns
            gnsInfo_new_dict[gns]["name"] = name
            gnsInfo_new_dict[gns]["creator"] = creator
            gnsInfo_new_dict[gns]["editor"] = editor
            gnsInfo_new_dict[gns]["create_time"] = create_time
            gnsInfo_new_dict[gns]["path"] = path
            gnsInfo_new_dict[gns]["type"] = gns_type
            gnsInfo_new_dict[gns]["rev"] = rev
            gnsInfo_new_dict[gns]["root"] = root
    return gnsInfo_new_dict
def insert_to_mogodb(conn,db_name, gnsInfo_dict,as7_json):
    # conn[db_name]["gnsInfo"].delete_many({})
    # conn.(db_name, "gnsInfo")
    db = conn[db_name]
    # db["gnsInfo"].insert_one({"a":2})
    info_dict = defaultdict(dict)
    gnsInfo_new_dict = get_gnsInfo_dict(gnsInfo_dict)
    for gns in gnsInfo_dict:
        gns_type = gnsInfo_dict[gns][-1]
        gnsInfo = gnsInfo_new_dict[gns]
        gns_pre = gns.split("//")[0] + "//"
        id_list = gns.split("//")[-1].split("/")
        if len(id_list) == 1:
            gnsInfo["child_file"] = []
            gnsInfo["child_folder"] = []
            info_dict[gns] = gnsInfo
            continue
        father_gns = gns_pre + "/".join(id_list[:-1])
        if int(gns_type) == 0:
            info_dict[father_gns]["child_file"].append(gns)
        else:
            gnsInfo["child_file"] = []
            gnsInfo["child_folder"] = []
            try:
                info_dict[father_gns]["child_folder"].append(gns)
            except Exception:
                pass
        info_dict[gns] = gnsInfo
    for gns, gnsInfo in info_dict.items():
        result = db["gnsInfo"].find_one(gnsInfo)
        if not result and gnsInfo:
            db["gnsInfo"].insert_one(gnsInfo)


class FileDownloadByGnsId():
    def __init__(self,entity_property_dict, relationship_property_dict,gnsid,conn,config_path,as7json,graph_name,words_path,templates_path,model_path,is_increment):
        self.gnsid = gnsid
        self.conn=conn
        self.entity_property_dict=entity_property_dict
        self.relationship_property_dict=relationship_property_dict
        self.config_path=config_path
        self.as7json=as7json
        self.ds_auth=as7json["ds_auth"]
        self.ds_address=as7json["ds_address"]
        self.ds_port=as7json["ds_port"]
        self.ds_name=as7json["dsname"]
        self.ds_id=as7json["id"]
        self.graph_name=graph_name
        self.mongo_gnsInfo_db = conn[graph_name]
        self.Config = self.get_config()
        self.gnsInfo_dict = defaultdict(dict)
        self.words_path = words_path
        self.templates_path = templates_path
        self.model_path = model_path
        self.is_increment = is_increment

    # def clean_mongo_collection(self):
    #     # m = Mongo()
    #     # 先清空collection
    #     for item in (list(self.entity_property_dict.keys())+list(self.relationship_property_dict.keys())):
    #         self.conn[self.graph_name][item].delete_many({})
            # self.conn as7_json["ds_name"], item)

    def get_config(self):
        config = ConfigParser()
        config.read(self.config_path, "utf-8")
        return config

    def get_all_gns(self):
        gns_list = []

        def _get_child_file(gns):
            if self.is_increment :
                coll = self.mongo_gnsInfo_db["gnsInfoChange"].find_one({"gns": gns})
            else :
                coll = self.mongo_gnsInfo_db["gnsInfo"].find_one({"gns": gns})
            gns_list.append(gns)
            try:
                gns_list.extend(coll["child_file"])
            except Exception:
                pass
            try:
                for folder_gns in coll["child_folder"]:
                    _get_child_file(folder_gns)
            except Exception:
                pass

        _get_child_file(self.gnsid)
        return gns_list

    def get_url(self, gns, name):
        ds_auth = self.as7json['ds_auth']
        print('开始获取AS token', __file__, 'get_url')
        ret_token, token_id = asToken.get_token(ds_auth)
        headers = {
            'Content-Type': 'application/json',
            'Authorization': 'Bearer {}'.format(token_id)
        }
        if "http" in self.ds_address:
            url = self.ds_address + ":" + str(self.ds_port) + self.Config['url_path'].get('fileDownloadPath')
            ds_address=self.ds_address.replace("https://","")
        else:
            url = "https://" + self.ds_address + ":" + str(self.ds_port) + self.Config['url_path'].get('fileDownloadPath')
        payload = json.dumps(
            {"docid": gns, "rev": "", "authtype": "", "reqhost": ds_address, "usehttps": False,
             "savename": name})
        response = requests.request("POST", url, headers=headers, data=payload, verify=False)
        return json.loads(response.text)

        # url=self.ds_address+":"+self.ds_port+self.Config['url_path'].get('fileDownloadPath')
        # payload = json.dumps(
        #     {"docid": gns})
        # response = requests.request("POST", url, headers=headers, data=payload, verify=False)
        # return json.loads(response.text)

    def get_file_content_by_url(self, authrequest):
        one_headers = {}
        for i in authrequest[2:]:
            one_headers[i[:i.index(":")].strip()] = i[i.index(":") + 1:].strip()
        response = requests.request(authrequest[0], authrequest[1], headers=one_headers,verify=False)
        return response

    def write_spo_with_response(self, response, file_info_dict):
        file_type = file_info_dict["name"].split(".")[-1]
        try:
            try:
                if file_type.lower() in ["pptx", "docx"]:
                    reader = io.BufferedReader(io.BytesIO(response.content))
                    Anyshare(self.conn,file_info_dict, self.graph_name,self.entity_property_dict,self.relationship_property_dict,self.as7json,file_type.lower(),self.words_path,self.templates_path,self.model_path)\
                        .process(reader)
                    del reader
                elif file_type.lower() in ["xlsx", "xls"]:
                    Anyshare(self.conn,file_info_dict, self.graph_name,self.entity_property_dict,self.relationship_property_dict,self.as7json,file_type.lower(),self.words_path,self.templates_path,self.model_path)\
                        .process(response.content)
                elif file_type.lower() in [ "doc","ppt","txt","pdf"]:
                    Anyshare(self.conn, file_info_dict, self.graph_name, self.entity_property_dict,
                             self.relationship_property_dict,self.as7json,file_type.lower(),self.words_path,self.templates_path,self.model_path).process(file_info_dict)
                else:
                    Anyshare(self.conn, file_info_dict, self.graph_name, self.entity_property_dict,
                             self.relationship_property_dict, self.as7json,file_type.lower(),self.words_path,self.templates_path,self.model_path).process("othertype")
            except Exception as e:
                print(repr(e))
                file_type="txt"
                Anyshare(self.conn, file_info_dict, self.graph_name, self.entity_property_dict,
                         self.relationship_property_dict, self.as7json, file_type.lower(),self.words_path,self.templates_path,self.model_path).process(file_info_dict)
            # elif file_type in ["txt"]:
            #     Anyshare(self.conn,file_info_dict, self.graph_name,self.entity_property_dict,self.relationship_property_dict)\
            #         .process(response.content.decode("gb18030"))

        except Exception as e:
            print(repr(e))

    def normalize_text(self,text):
        text = re.sub(r"[\n\t\'\"]", " ", text)
        text = text.replace("\\", "\\\\").strip()
        return text

    def get_path_property_dict(self):
        gns_list = self.get_all_gns()
        for i, gns in enumerate(gns_list):
            if self.is_increment:
                coll = self.mongo_gnsInfo_db["gnsInfoChange"].find_one({"gns": gns})  # gns具有唯一性，故用find_one
            else:
                coll = self.mongo_gnsInfo_db["gnsInfo"].find_one({"gns": gns})  # gns具有唯一性，故用find_one
            fileInfo_dict = dict(coll)
            self.normalize_text(fileInfo_dict["path"])
            # self.path_property_dict[path] = fileInfo_dict

    def get_attribute_from_es(self, gns, address,port, esapi):
        import requests
        # address = address.replace("https", "http")
        esgnsid = gns.split('/')[-1]
        url = esapi.format(address,port, esgnsid)
        payload = {}
        headers = {}
        response = requests.request("GET", url, headers=headers, data=payload,verify=False)
        res = response.json()
        _source = res["_source"]
        return _source

    def start(self):
        def _get_file_content():
            try:
                url_response = self.get_url(gns, name)
                response = self.get_file_content_by_url(url_response["authrequest"])
                if response.status_code == 200:
                    ds_address = self.as7_json["ds_address"]
                    version = otl_dao.getversion(ds_address)
                    retcode ,source = otl_dao.getinfofromas(self.as7_json,["_source"],gns,version)
                    if retcode == 500:
                        print("get file content error: {}".format(name))
                    else:
                        fileInfo_dict["ds_address"] = self.ds_address
                        fileInfo_dict["ds_id"] = self.ds_id
                        fileInfo_dict["file_type"] = source["ext"]
                        fileInfo_dict["editor"] = source["editor"]
                        fileInfo_dict["creator"] = source["creator"]
                        fileInfo_dict["create_time"] = time.strftime("%Y-%m-%d %H:%M:%S",time.localtime(int(str(source["created"])[0:10])))
                        fileInfo_dict["modified_time"]  = time.strftime("%Y-%m-%d %H:%M:%S",time.localtime(int(str(source["modified"])[0:10])))
                        fileInfo_dict["label"] = source["tags"]
                        fileInfo_dict["rev"] = source["rev"]
                        if "content" in source:
                            fileInfo_dict["content"] = source["content"]
                        self.write_spo_with_response(response, fileInfo_dict)
                else:
                    print("get file content error: {}".format(name))
                    # raise IOError
            except Exception:
                print("get file content error: {}".format(name))
                # raise IOError

        # mongo 获取所有有效文档及文件夹的信息
        gns_list = self.get_all_gns()
        # 导入mongo
        # self.clean_mongo_collection()

        for i, gns in enumerate(gns_list):
            if self.is_increment:

                coll = self.mongo_gnsInfo_db["gnsInfoChange"].find_one({"gns": gns})
                if coll :
                    fileInfo_dict = dict(coll)
                    name = fileInfo_dict["name"]
                    gns_type = fileInfo_dict["type"]
                    change_type = fileInfo_dict["change_type"]
                    if change_type in [1, 2, 3]:
                        if int(gns_type) == 0:
                            # if i >= 8065:
                            _get_file_content()
            else:
                coll = self.mongo_gnsInfo_db["gnsInfo"].find_one({"gns": gns})  # gns具有唯一性，故用find_one
                fileInfo_dict = dict(coll)
                name = fileInfo_dict["name"]
                gns_type = fileInfo_dict["type"]
                if int(gns_type) == 0 :
                    # if i >= 8065:
                    _get_file_content()




class ImportOrientFromMongo():
    def __init__(self, graphname, graph_json, conn):
        self.conn=conn
        self.mongo_db = self.conn[graphname]
        self.baseInfo_json = graph_json["graph_baseInfo"][0]
        self.address = self.baseInfo_json["graphDBAddress"]  # 图数据库ip
        self.graph_db = self.baseInfo_json["graph_DBName"]

    # def get_mongo_data(self):
    #     mongo_spo_data = {}
    #     for item in (list(entity_property_dict.keys()) + list(relationship_property_dict.keys())):
    #         mongo_spo_data[item] = list(self.mongo_db[item].find())
    #     return mongo_spo_data

    def oriendb_http(self, sql):
        import requests
        from requests.auth import HTTPBasicAuth
        orient_url = "http://{}:{}/command/{}/sql".format(self.address, "2480", self.graph_db)
        body = {"command": sql}
        or_res = requests.post(url=orient_url, json=body, auth=HTTPBasicAuth("admin", "admin"))
        if or_res.status_code == 200:
            print("-------Sucess------")
        return or_res.json()

    def import2orient(self):
        from dao.task_dao import task_dao
        ret = task_dao.getGraphDBbyIp(self.address)
        rec_dict = ret.to_dict('records')
        rec_dict = rec_dict[0]
        username = rec_dict["db_user"]
        password = rec_dict["db_ps"]
        password = commonutil.DecryptBybase64(password)
        graph_DBPort = rec_dict["port"]
        # # 删除数据库
        # self.oriendb_http("drop database %s"%(self.graph_db))
        # 查询是否存在，不存在在创建
        url = 'http://' + self.address + ':' + str(graph_DBPort) + '/database/' + self.graph_db
        r = requests.get(url, auth=(username, password))
        if r.status_code != 200:
            requests.post(
                'http://' + self.address + ':' + str(graph_DBPort) + '/database/' + self.graph_db + '/plocal',
                auth=(username, password))
            # url = 'http://' + self.address + ':' + str(graph_DBPort) + '/database/' + self.graph_db + '/plocal'

        # mongo_spo_data = self.get_mongo_data()
        # Mongo2orient(self.graph_db, self.oriendb_http, mongo_spo_data, entity_property_dict,relationship_property_dict)\
        #     .import2orient()

class Anyshare():
    def __init__(self, conn,file_info_dict,graph_name,entity_property_dict,relationship_property_dict,as7json,file_type,words_path,templates_path,model_path):
        self.graph_name=graph_name
        self.file_info_dict = file_info_dict
        self.conn=conn
        self.as7json=as7json
        self.address = as7json["ds_address"]
        self.ds_port = as7json["ds_port"]
        self.ds_id = as7json["id"]
        self.file_name = self.normalize_text(file_info_dict["name"])
        self.file_path = self.normalize_text(file_info_dict["path"])
        self.file_gns = file_info_dict["gns"]
        self.file_type =file_type
        self.text_list = []
        self.level_list = []
        self.head_list = []
        self.head2text_spo_list = []
        self.document2head_spo_list = []
        self.document2text_spo_list = []
        self.label2document_spo_list = []
        self.text2text_spo_list = []
        self.desc2document_spo_list = []
        self.desc2label_spo_list = []
        self.entity_property_dict = entity_property_dict
        self.relationship_property_dict = relationship_property_dict
        self.sentence_list = []
        self.sentence_frag_list = []
        self.spo_list = []
        self.words_path = words_path
        self.templates_path = templates_path
        self.model_path=model_path

    def split_text_and_build_spo(self, text, s, s_type, s_pro, spo_list):
        num = len(text) // 1000
        text_list = []
        if num >= 1:
            for i in range(1, num + 1):
                now_text = text[(i - 1) * 1000:i * 1000]
                text_list.append(now_text)
                if len(text_list) > 1:
                    pre_text = text_list[-2]
                    item_dict = {"s": pre_text, "p": "text2text", "o": now_text, "o_pro": {}, "p_pro": {},
                                 "s_pro": {}, "path": self.file_path}
                    self.text2text_spo_list.append(item_dict)
        else:
            text_list.append(text)
        for text in text_list:
            if s_type == "head":
                item_dict = {"s": s, "p": "head2text", "o": text, "o_pro": {}, "p_pro": {},
                             "s_pro": s_pro, "path": self.file_path}
            else:
                item_dict = {"s": s, "p": "document2text", "o": text}
            spo_list.append(item_dict)
    def get_attribute_from_es(self,gns,address,port,esapi):
        import requests
        # address = address.replace("https", "http")
        esgnsid = gns.split('/')[-1]
        # address = "http://10.4.69.44"
        # esgnsid = "1D2666C81D6B4A6B8999DD13082F28A2"  # json 515D42AD63FF424B9DEFA4F40270F724
        url = esapi.format(address,port, esgnsid)
        payload = {}
        headers = {}
        response = requests.request("GET", url, headers=headers, data=payload)
        res = response.json()
        _source = res["_source"]
        return _source

    def process(self, reader):
        if self.file_type.lower() in ["docx"]:
            self._docx_processor(reader)
        # elif self.file_type == "doc":
        #     self._doc_processor(reader)
        # elif self.file_type == "pdf":
        #     self._pdf_processor(reader)
        elif self.file_type.lower() in ["pptx"]:
            self._pptx_processor(reader)
        elif self.file_type.lower() in ["ppt","doc","txt","pdf"]:##走es接口
            if "content" in reader:
                self._es_processor(reader)
        elif self.file_type.lower() in ["xlsx", "xls"]:
            self._excel_processor(reader)
        self.get_label2document()

        # print("start extracting label description ")
        import datetime
        # descs = datetime.datetime.now()
        # print(descs)
        # ede_model = EntityDescribeExtractByRoleAnalysis(self.model_path)

        # es = extraction_start(self.model_path)
        # self.get_sentence_list()
        # self.label_dict_list = ede_model.extract_info(self.sentence_list)
        # desce = datetime.datetime.now()
        # print("description spend {}s".format((desce-descs).total_seconds()))
        # # self.ltp ,self.spo_list = es.run(self.sentence_list)
        # # self.label_dict = Recognizor(self.ltp,self.spo_list,self.words_path,self.templates_path,self.model_path).filter_by_p_words()
        # # self.label_dict = self.filter_entity_describe(self.ltp,self.label_dict)
        # self.get_spo_lists(self.label_dict_list)
        # 保存三元组
        print("spo2mongo start ")
        mongos = datetime.datetime.now()
        print(mongos)
        self.spo2mongo()
        print("spo2mongo finished ")
        mongoe = datetime.datetime.now()
        print(mongoe)
        print("spo2mongo spends {}s ".format((mongoe-mongos).total_seconds()))



    # def filter_entity_describe(self,ltp,input_label_dict):
    #     entity_describe_dict = {}
    #     for key, value in input_label_dict.items():
    #         sentence = key
    #         entity_describe_dict[sentence] =[]
    #         seg, hidden = ltp.seg([sentence])
    #         words = seg[0]
    #         pos = ltp.pos(hidden)[0]
    #         assert len(words) == len(pos)
    #         for s_o_str in value:
    #             sub, obj = s_o_str[0],s_o_str[1]
    #             # sub_list = re.split("（", sub)
    #             # sub = sub_list[0]
    #             try:
    #                 sub_ind = sentence.index(sub)
    #             except Exception as e:
    #                 continue
    #             sub_end = sub_ind + len(sub)
    #
    #             cal_i = 0
    #             word_seq = []
    #             for i, w in enumerate(words):
    #                 if cal_i >= sub_ind and cal_i < sub_end:
    #                     word_seq.append((i, w))
    #                 cal_i += len(w)
    #
    #             last_ind = word_seq[-1][0]
    #             if pos[last_ind] not in ["n", "nz", "i"]:
    #                 continue
    #             if pos[word_seq[0][0]] in ["r"]:
    #                 continue
    #             entity_describe_dict[sentence].append([sub, obj])
    #     return entity_describe_dict
    def get_spo_lists(self,label_dict_list):
        db = self.conn[self.graph_name]
        if label_dict_list:
            for label_desc in label_dict_list:
                sentence = label_desc["sentence"]
                label = label_desc["entity"]
                label_dict_first = {"name": label, "adlabel_kcid": label, "type_nw":"true"}
                ss = db["label"].find_one(label_dict_first)
                if ss:
                    label2document_item_dict = {"s": label, "p": "label2document", "o": self.file_path,"type":"type_sa"}
                    self.label2document_spo_list.append(label2document_item_dict)
                    desc2label_item_dict = {"s": sentence, "p": "desc2label", "o": label}
                    self.desc2label_spo_list.append(desc2label_item_dict)
                    desc2document_item_dict = {"s": sentence, "p": "desc2document", "o":label }
                    self.desc2document_spo_list.append(desc2document_item_dict)



    def get_sentence_list(self,content):
        text = content.replace("\n", "")
        self.sentence_list = [sentence.strip() for sentence in re.split(r"[。\n]", text) if len(sentence) > 2]

    def get_label2document(self):
        for l in self.file_info_dict["label"]:
            item_dict = {"s": l , "p": "label2document", "o": self.file_name, "type":"type_as"}
            self.label2document_spo_list.append(item_dict)
    def _es_processor(self,reader):
        content = reader["content"]
        if len(content) > 1:
            self.split_text_and_build_spo(self.normalize_text(content), self.file_name, "document", {},self.document2text_spo_list)

    def _docx_processor(self, reader):
        document = docx.Document(reader)
        parent_head = ()
        self.head_list.append((self.file_name, 0))
        self.level_list.append(0)
        text_list = []
        for p in document.paragraphs:
            text_list.append(p.text)
            if p.style.name.startswith('Heading'):
                p_level = int(p.style.name.split(" ")[-1])
                head = self.normalize_text(p.text)
                self.level_list.append(p_level)
                self.head_list.append((head, p_level))
                parent_head = (head, p_level)
                item_dict = {"s": self.file_name, "p": "document2chapter", "o": head, "o_pro": {"level": p_level}}
                self.document2head_spo_list.append(item_dict)
            elif p.style.name == 'Normal' or p.style.name == '云正':
                text = self.normalize_text(p.text)
                if len(text) > 1 and parent_head:
                    item_dict = {"s": parent_head[0], "p": "chapter2text", "o": text,
                                 "s_pro": {"level": parent_head[1]}}
                    self.head2text_spo_list.append(item_dict)
                if len(text) > 1 and not parent_head:
                    item_dict = {"s": self.file_name, "p": "document2text", "o": text}
                    self.document2text_spo_list.append(item_dict)

        if not self.head2text_spo_list and text_list:
            text = self.normalize_text(" ".join(text_list))
            self.split_text_and_build_spo(text, self.file_name, "document", {},self.document2text_spo_list)

    # def _doc_processor(self, reader):
    #     # tika.initVM()
    #     # from tika import parser
    #     parsed = parser.from_buffer(reader)
    #     content = parsed["content"]
    #     text_list = [sentence.strip() for sentence in content.split("\n") if sentence]
    #     # 因为doc读取后的文档无格式，排版也比较乱，故不提取结构
    #     text = self.normalize_text(" ".join(text_list))
    #     item_dict = {"s": self.file_name, "p": "document2text", "o": text, "o_pro": {}, "p_pro": {},
    #                  "s_pro": {}, "path": self.file_path}
    #     self.document2text_spo_list.append(item_dict)

    def _pdf_processor(self, reader):
        # pdf里面有表格需要单独处理
        def ___head_extract(line):
            if line:
                if "....." in line:
                    head_add_num = line.split(".....")[0]
                    num_re = re.findall(r"[0-9.]{2,}|[第一二三四五六七八九十0-9篇章节：:]{2,}", head_add_num)
                    num = num_re[0].strip() if num_re else ""
                    head = head_add_num.split(num)[-1].strip() if num else head_add_num
                    c = num.count(".")
                    head_level_dict[c + 1].append(head)
                    self.level_list.append(c + 1)
                    item_dict = {"s": self.file_name, "p": "document2chapter", "o": head, "o_pro": {"level": c + 1}}
                    self.document2head_spo_list.append(item_dict)
                    head_reg_list.append((num, head, head_add_num, c + 1))

        head_list, text_list = [], []
        pdf = pdfplumber.load(reader)
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_list.append(page_text)
        try:
            head_reg_list = []
            head_level_dict = {key: [] for key in range(6)}
            # 锁定目录部分
            catalog_text = ""
            for text in text_list:
                if re.findall(r"[.]{5,}[\s]*[1-9]", text):
                    catalog_text = catalog_text + "\n" + text
            for line in catalog_text.split("\n"):
                ___head_extract(line)
            # 提取正文部分
            sentence_list = "\n".join(text_list).split("\n")
            start_index = 0
            flag = 1
            now_head_reg = ""
            text_parent_head_list = []
            for i, sentence in enumerate(sentence_list):
                sentence_drop_space = sentence.replace(" ", "")
                if not head_reg_list:
                    last_sentence = "".join(sentence_list[start_index + 1:]) if i > start_index else ""
                    if now_head_reg[1] in last_sentence:
                        text1 = last_sentence.split(now_head_reg[1])[0]
                        text1 = "".join(text1.split(now_head_reg[0])[:-1])
                        text2 = last_sentence.split(now_head_reg[1])[1]
                    # text = "".join(sentence_list[start_index + 1:i]) if i > start_index else ""
                    if text1:
                        text1 = self.normalize_text(text1)
                        head_mix = text_parent_head_list[-2]
                        if head_mix[0]:
                            self.split_text_and_build_spo(text1, self.normalize_text(head_mix[0]), "head",
                                                          {"level": head_mix[1]},
                                                          self.head2text_spo_list)
                        else:
                            self.split_text_and_build_spo(text1, self.file_name, "document", {},
                                                          self.document2text_spo_list)
                    if text2:
                        text2 = self.normalize_text(text2)
                        if head_mix[0]:
                            self.split_text_and_build_spo(text2, self.normalize_text(now_head_reg[1]), "head",
                                                          {"level": now_head_reg[-1]},
                                                          self.head2text_spo_list)
                        else:
                            self.split_text_and_build_spo(text2, self.file_name, "document", {},
                                                          self.document2text_spo_list)

                    break
                if flag:
                    now_head_reg = head_reg_list.pop(0)
                    now_head_level = now_head_reg[-1]
                    text_parent_head_list.append((now_head_reg[1], now_head_level))
                    flag = 0
                if now_head_reg[2].replace(" ", "") in sentence_drop_space and "." * 5 not in sentence_drop_space:
                    if start_index:
                        text = "".join(sentence_list[start_index + 1:i]) if i > start_index else ""
                        if text:
                            text = self.normalize_text(text)
                            head_mix = text_parent_head_list[-2]
                            if head_mix:
                                self.split_text_and_build_spo(text, self.normalize_text(head_mix[0]), "head",
                                                              {"level": head_mix[1]},
                                                              self.head2text_spo_list)
                            else:
                                self.split_text_and_build_spo(text, self.file_name, "document", {},
                                                              self.document2text_spo_list)
                    flag = 1
                    start_index = i
            self._normalize_level()
        except Exception:
            pass
        # print(self.head2text_spo_list)
        if not self.head2text_spo_list and text_list:
            text = self.normalize_text(" ".join(text_list))
            self.split_text_and_build_spo(text, self.file_name, "document", {},self.document2text_spo_list)


    def _pptx_processor(self, reader):
        def __iter_shape(shape, head, shape_head_list):
            if type(shape) == pptx.shapes.group.GroupShape:
                for sshape in shape.shapes:
                    __iter_shape(sshape, head, shape_head_list)
            else:
                if shape.has_text_frame or shape.has_table:
                    shape_head_list.append((shape, head))

        def __combine_content(content, head, head_content_dict):
            try:
                head_content_dict[head] = head_content_dict[head] + " " + content
            except Exception:
                head_content_dict[head] = content
            return head_content_dict

        def __pptx_paragraph_text(ppt):
            def find_head(shape):
                if type(shape) == pptx.shapes.group.GroupShape:
                    for sshape in shape.shapes:
                        target_head=find_head(sshape)
                        return target_head
                else:
                    if shape.has_text_frame:
                        target_head=shape.text_frame.text
                        result=re.findall(r"，。、,.;；",target_head)
                        if result:
                            return None
                        else:
                            return target_head
            shape_head_list, head_list = [], []
            head_list.append(self.file_name)
            slide_list = [slide for slide in ppt.slides]
            for slide in slide_list:
                shapes = slide.shapes
                if slide.shapes.title:
                    head = slide.shapes.title.text
                else:
                    for shape in shapes:
                        head = find_head(shape)
                        break
                if not head:
                    head = head_list[-1]
                head_list.append(head)
                for shape in shapes:
                    __iter_shape(shape, head, shape_head_list)
            paragraphs = [(paragraph, head) for shape, head in shape_head_list
                          if shape.has_text_frame for paragraph in shape.text_frame.paragraphs]
            item_list = [(p.text, head) for p, head in paragraphs if p.text and p.text != head]
            for shape, head in shape_head_list:
                if shape.has_table:
                    cell_text_list = [cell.text for cell in shape.table.iter_cells() if cell.text_frame]
                    item_list.append((" ".join(cell_text_list), head))

            head_content_dict = defaultdict()
            for content, head in item_list:
                head_content_dict = __combine_content(content, head, head_content_dict)
            return head_content_dict

        ppt = pptx.Presentation(reader)
        head_content_dict = __pptx_paragraph_text(ppt)
        for head, content in head_content_dict.items():
            item_dict = {"s": self.file_name, "p": "document2chapter", "o": self.normalize_text(head),
                         "o_pro": {"level": 1}}
            self.document2head_spo_list.append(item_dict)
            self.split_text_and_build_spo(self.normalize_text(content), self.normalize_text(head), "head", {"level": 1},
                                          self.head2text_spo_list)
            # self.head2text_spo_list.append(item_dict)
        # print(self.head2text_spo_list)

    def _excel_processor(self, reader):
        excel = xlrd.open_workbook(file_contents=reader)
        sheet_names = excel.sheet_names()
        for sheet in sheet_names:
            table = excel.sheet_by_name(sheet)
            item_list = [table.cell_value(row, col) for row, col in product(range(table.nrows), range(table.ncols))]
            item_dict = {"s": self.file_name, "p": "document2chapter", "o": self.normalize_text(sheet),
                         "o_pro": {"level": 1}}
            self.document2head_spo_list.append(item_dict)
            if item_list:
                text = self.normalize_text(" ".join(item_list))
                self.split_text_and_build_spo(text, self.normalize_text(sheet), "head", {"level": 1},
                                              self.head2text_spo_list)

    # def _txt_processor(self, reader):
    #     item_dict = {"s": self.file_name, "p": "document2text", "o": self.normalize_text(reader), "o_pro": {},
    #                  "p_pro": {},
    #                  "s_pro": {}, "path": self.file_path}
    #     self.document2text_spo_list.append(item_dict)

    def spo2mongo(self):
        db = self.conn[self.graph_name]
        document_property_list = [item[0] for item in self.entity_property_dict["document"]]
        head_property_list = [item[0] for item in self.entity_property_dict["chapter"]]
        text_property_list = [item[0] for item in self.entity_property_dict["text"]]
        folder_property_list = [item[0] for item in self.entity_property_dict["folder"]]
        label_property_list = [item[0] for item in self.entity_property_dict["label"]]
        desc_property_list = [item[0] for item in self.entity_property_dict["desc"]]
        # folder2folder_property_list = [item[0] for item in self.relationship_property_dict["folder2folder"]]
        # folder2document_property_list = [item[0] for item in self.relationship_property_dict["folder2document"]]
        # document2head_property_list = [item[0] for item in self.relationship_property_dict["document2chapter"]]
        # document2text_property_list = [item[0] for item in self.relationship_property_dict["document2text"]]
        # head2text_property_list = [item[0] for item in self.relationship_property_dict["chapter2text"]]
        # 导入document
        property_dict = {document_property_list[0]: self.normalize_text(self.file_name),###name
                         document_property_list[1]: self.normalize_text(self.file_path),###path
                         document_property_list[2]: self.file_info_dict["creator"],###creator
                         document_property_list[3]: self.file_info_dict["create_time"],###create_time
                         document_property_list[4]: self.file_info_dict["editor"],###editor
                         document_property_list[5]: self.normalize_text(self.file_info_dict["gns"]),###gns
                         document_property_list[6]: self.file_info_dict["file_type"],
                         document_property_list[7]: self.file_info_dict["modified_time"],
                         document_property_list[8]: self.file_info_dict["rev"]
                         }
        # if len(self.file_info_dict["label"])>0:
        #     for l in self.file_info_dict["label"]:
        #         label_dict = {
        #             label_property_list[0]: l
        #         }
        #         db["label"].insert_one(label_dict)
        property_dict["ds_id"] = self.file_info_dict["ds_id"]
        db["document"].insert_one(property_dict)

        # 导入folder
        path_folder = self.file_path.replace(self.file_name, "")
        folder_list = path_folder.split("/")[:-1]
        gns_pre = "gns://"
        gns_list = self.file_info_dict["gns"].replace(gns_pre, "").split("/")[:-1]
        assert len(gns_list) == len(folder_list)
        done_head_list, done_text_list = [], []
        for i, folder in enumerate(folder_list):
            folder_path = self.normalize_text("/".join(folder_list[:i + 1]))
            gns = self.normalize_text(gns_pre + "/".join(gns_list[:i + 1]))
            try:
                ds_address = self.as7_json["ds_address"]
                version = otl_dao.getversion(ds_address)
                ret_code, source = otl_dao.getinfofromas(self.as7_json, ["_source"], gns, version)
                # ret_code ,source=otl_dao.getinfofromas(self.ds_id,["created"], gns)
                # source=self.get_attribute_from_es(gns,self.address,self.ds_port,self.esapi)
                create_time=time.strftime("%Y-%m-%d %H:%M:%S",time.localtime(int(str(source["created"])[0:10])))
                rev = source["rev"]
            except Exception:
                try:
                    entry_code, entry_obj = otl_dao.get_entry_doc(self.as7json)
                    create_time = entry_obj[gns]["created_at"]
                    from dateutil import parser
                    import datetime
                    from pytz import timezone
                    def convert(rfc3339):
                        converted_date = datetime.datetime.fromtimestamp(parser.parse(rfc3339).timestamp()).strftime('%Y-%m-%d %H:%M:%S')
                        return converted_date
                    create_time=convert(create_time)
                    rev = entry_obj[gns]["rev"]
                except Exception:
                    create_time=""
                    rev = ""
            property_dict = {folder_property_list[0]: folder,
                             folder_property_list[1]: folder_path,
                             folder_property_list[2]: gns,
                             folder_property_list[3]: create_time,
                             folder_property_list[4]: rev
                             }
            result = db["folder"].find_one(property_dict)
            if not result:
                property_dict["ds_id"] =str( self.file_info_dict["ds_id"])
                db["folder"].insert_one(property_dict)


        root = folder_list.pop(0)
        root_gns = gns_pre + gns_list.pop(0)
        latest_folder_gns = root_gns
        latest_folder = root
        for i, folder in enumerate(folder_list):
            folder_gns = root_gns + "/" + "/".join(gns_list[:i + 1])
            # folder2document_property_list属性赋值
            item_dict = {"s": latest_folder,
                         "p": "folder2folder",
                         "o": folder,
                         "s_pro": {"gns": latest_folder_gns},
                         "p_pro": {"name": "folder2folder"},
                         "o_pro": {"gns": folder_gns}}
            result = db["folder2folder"].find_one(item_dict)
            if not result:
                item_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                db["folder2folder"].insert_one(item_dict)
            latest_folder_gns = folder_gns
            latest_folder = folder
        item_dict = {"s": latest_folder,
                     "p": "folder2document",
                     "o": self.file_name,
                     "s_pro": {"gns": latest_folder_gns},
                     "p_pro": {"name": "folder2document"},
                     "o_pro": {"gns": self.file_gns}}
        item_dict["ds_id"] =str(self.file_info_dict["ds_id"])
        db["folder2document"].insert_one(item_dict)

        if self.head2text_spo_list:
            for item in self.head2text_spo_list:
                s, p, o, s_pro, = item["s"], item["p"], item["o"], item["s_pro"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)

                head_property_dict = {head_property_list[0]: self.normalize_text(s),head_property_list[1]:self.normalize_text(self.file_path), head_property_list[2]:  s_pro["level"] }
                head_property_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                db["chapter"].insert_one(head_property_dict)
                text_property_dict = {text_property_list[0]: o}
                text_property_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                db["text"].insert_one(text_property_dict)


                item_dict = {"s": s,
                             "p": "chapter2text",
                             "o": o,
                             "s_pro": head_property_dict,
                             "p_pro": {"name": "chapter2text"},
                             "o_pro": text_property_dict}
                item_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                db["chapter2text"].insert_one(item_dict)

        if self.document2head_spo_list:
            for item in self.document2head_spo_list:
                s, p, o, =  item["s"], item["p"], item["o"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)
                property_dict = {head_property_list[0]: o,
                                 head_property_list[1]: self.normalize_text(self.file_path),
                                 head_property_list[2]: item["o_pro"]["level"]
                                 }
                if s not in done_head_list:
                    property_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                    db["chapter"].insert_one(property_dict)
                    done_head_list.append(s)
                item_dict = {"s": s,
                             "p": "document2chapter",
                             "o": o,
                             "s_pro": {"gns": self.file_gns},
                             "p_pro": {"name": "document2chapter"},
                             "o_pro": property_dict}
                db["document2chapter"].insert_one(item_dict)
        if self.label2document_spo_list:
            for item in self.label2document_spo_list:
                s, p, o, db_type = item["s"],item["p"],item["o"],item["type"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)
                # name #adlabel_kcid #kc_topic_tags #type_as #type_sa #type_nw #type_kc
                property_dict = {label_property_list[0]: s }
                property_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                property_dict['adlabel_kcid'] = s
                commonutil.updatemongo(db,"label",property_dict,db_type)
                item_dict = {"s": s,
                             "p": "label2document",
                             "o": o,
                             "s_pro": property_dict,
                             "p_pro": {"name": "label2document"},
                             "o_pro": {"gns": self.file_gns}}
                db["label2document"].insert_one(item_dict)


        if self.document2text_spo_list:
            for item in self.document2text_spo_list:
                s, p, o = item["s"],item["p"],item["o"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)

                property_dict = {text_property_list[0]: o }
                property_dict["ds_id"] = str(self.file_info_dict["ds_id"])


                db["text"].insert_one(property_dict)

                item_dict = {"s": s,
                             "p": "document2text",
                             "o": o,
                             "s_pro": {"gns": self.file_gns},
                             "p_pro": {"name": "document2text"},
                             "o_pro": property_dict}
                db["document2text"].insert_one(item_dict)

        if self.text2text_spo_list:
            for item in self.text2text_spo_list:
                s, p, o = item["s"], item["p"], item["o"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)
                property_dict1 = {text_property_list[0]: s}
                property_dict2 = {text_property_list[0]: o}
                item_dict = {"s": s,
                             "p": "text2text",
                             "o": o,
                             "s_pro": property_dict1,
                             "p_pro": {"name": "text2text"},
                             "o_pro": property_dict2}
                db["text2text"].insert_one(item_dict)



        if self.desc2label_spo_list:
            for item in self.desc2label_spo_list:
                s, p, o = item["s"], item["p"], item["o"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)
                desc_property_dict = {desc_property_list[0]: s}
                desc_property_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                db["desc"].insert_one(desc_property_dict)
                property_dict2 = {label_property_list[0]: o}
                property_dict2['adlabel_kcid'] = o
                item_dict = {"s": s,
                             "p": "desc2label",
                             "o": o,
                             "s_pro": desc_property_dict,
                             "p_pro": {"name": "desc2label"},
                             "o_pro": property_dict2}
                db["desc2label"].insert_one(item_dict)
        if self.desc2document_spo_list:
            for item in self.desc2document_spo_list:
                s, p, o = item["s"], item["p"], item["o"]
                s = self.normalize_text(s)
                o = self.normalize_text(o)
                desc_property_dict = {desc_property_list[0]: s}
                desc_property_dict["ds_id"] = str(self.file_info_dict["ds_id"])
                item_dict = {"s": s,
                             "p": "desc2document",
                             "o": o,
                             "s_pro": desc_property_dict,
                             "p_pro": {"name": "desc2document"},
                             "o_pro": {"gns": self.file_gns}}
                db["desc2document"].insert_one(item_dict)



    def normalize_text(self, text):
        result = re.findall(r"[0-9]+[.][0-9\s] *", text)
        if result:
            if text.startswith(result[0]):
                text = text.replace(result[0], "")
        text = re.sub(r"[\n\t\'\"]", " ", text).strip()
        return text


    # 规范化标题等级排序
    def _normalize_level(self):
        self.level_list = list(set(self.level_list))
        self.level_list.sort()
        head2text_spo_list_copy, document2head_spo_list_copy = [], []
        for i in range(len(self.head2text_spo_list)):
            one, two, three, four = self.head2text_spo_list[i]
            head2text_spo_list_copy.append((one, two, three, self.level_list.index(four)))
        for i in range(len(self.document2head_spo_list)):
            one, two, three, four = self.document2head_spo_list[i]
            document2head_spo_list_copy.append((one, two, three, self.level_list.index(four)))
        self.head2text_spo_list = head2text_spo_list_copy
        self.document2head_spo_list = document2head_spo_list_copy
class Mongo2orient():
    def __init__(self, sql_command, spo_data, entity_property_dict, relationship_property_dict):
        self.entity_property_dict = entity_property_dict
        self.relationship_property_dict = relationship_property_dict
        self.sql_command = sql_command
        self.spo_data = spo_data

    def normalize_text(self, text):
        text = re.sub(r"[\n\t\'\"]", " ", text)
        text = text.replace("\\", "\\\\").strip()
        return text

    # def create_ontology(self, ):
        # 文档实体及属性
        # for entity in self.entity_property_dict:
        #     create_entity_sql = "create class {} extends V".format(entity)
        #     try:
        #         self.sql_command(create_entity_sql)
        #         for pro, pro_type in self.entity_property_dict[entity]:
        #             create_entity_property_sql = "create property {}.`{}` {}".format(entity, pro, pro_type)
        #             self.sql_command(create_entity_property_sql)
        #     except:
        #         pass
        # # 创建关系类
        # for relationship in self.relationship_property_dict:
        #     create_relationship_sql = "create class {} extends E".format(relationship)
        #     try:
        #         self.sql_command(create_relationship_sql)
        #         for pro, pro_type in self.relationship_property_dict[relationship]:
        #             create_relationship_property_sql = "create property {}.`{}` {}".format(relationship, pro,
        #                                                                                    pro_type)
        #             self.sql_command(create_relationship_property_sql)
        #     except:
        #         pass

    # def create_entity(self):
    #
    #     def __command_create(entity, pro_dict):
    #         create_vertex_sql = "create vertex {}".format(entity)
    #         num = 0
    #         if pro_dict:
    #             for pro in pro_dict:
    #                 pro_value = self.normalize_text(str(pro_dict[pro]))
    #                 if num == 0:
    #                     create_vertex_sql += " set `{}`='{}'".format(pro, pro_value)
    #                 else:
    #                     create_vertex_sql += ", `{}`='{}'".format(pro, pro_value)
    #                 num += 1
    #         try:
    #             self.sql_command(create_vertex_sql)
    #         except Exception as e:
    #             print(e)
    #
    #     for entity in self.entity_property_dict.keys():
    #         for item in self.spo_data[entity]:
    #             if "_id" in item:
    #                 del item['_id']
    #             print(item)
    #             __command_create(entity, item)

    def create_index(self):
        start = time.clock()
        for entity in self.entity_property_dict:
            # orient.execute("DROP INDEX {}.property_index".format(entity))
            if entity == "chapter":
                index_property_list = [item[0] for item in self.entity_property_dict[entity]]
                create_property_index_sql = "CREATE INDEX {}.property_index on {} ({}) UNIQUE" \
                    .format(entity, entity, ",".join(index_property_list))
            else:
                index_property = [item[0] for item in self.entity_property_dict[entity]][1]  # 其他有一个即可确定唯一
                create_property_index_sql = "CREATE INDEX {}.property_index on {} ({}) UNIQUE" \
                    .format(entity, entity, index_property)
            try:
                self.sql_command(create_property_index_sql)
            except Exception:
                pass
        for relationship in self.relationship_property_dict:
            property_list = [item[0] for item in self.relationship_property_dict[relationship]]
            create_property_index_sql = "CREATE INDEX {}.property_index on {} ({}) UNIQUE" \
                .format(relationship, relationship, ",".join(property_list))
            try:
                self.sql_command(create_property_index_sql)
            except Exception:
                pass
        print("create index use:", (time.clock() - start))

    def create_relationship(self, ):
        def __command_create(pro_class, start_sql, end_sql, pro_dict):
            relation_create_sql = "create edge {} from ({}) to ({})" \
                .format(pro_class, start_sql, end_sql)
            num = 0
            if pro_dict:
                for pro in pro_dict:
                    pro_value=self.normalize_text(str(pro_dict[pro]))
                    if num == 0:
                        relation_create_sql += " set `{}`='{}'".format(pro, pro_value)
                    else:
                        relation_create_sql += ", `{}`='{}'".format(pro, pro_value)
                    num += 1
            try:
                self.sql_command(relation_create_sql)
            except Exception as e:
                print(e)

        def _get_select_sql(obj, property_dict):
            sql = 'select from `{}` where '.format(obj)
            flag = 0
            for i, (property_key, property_value) in enumerate(property_dict.items()):
                property_value=self.normalize_text(str(property_value))
                if property_value in ["None", "null", " "]:
                    flag += 1
                    continue
                if i == 0:
                    sql += "`{}`='{}'".format(property_key, property_value)
                elif flag == i:
                    sql += "`{}`='{}'".format(property_key, property_value)
                else:
                    sql += " and `{}`='{}'".format(property_key, property_value)
            return sql

        for relationship in self.relationship_property_dict.keys():
            s_class = relationship.split("2")[0]
            o_class = relationship.split("2")[-1]
            for item in self.spo_data[relationship]:
                s, p, o, s_pro, p_pro, o_pro = \
                    item["s"], item["p"], item["o"], item["s_pro"], item["p_pro"], item["o_pro"]
                if "_id" in s_pro:
                    del s_pro["_id"]
                if "_id" in p_pro:
                    del p_pro["_id"]
                if "_id" in o_pro:
                    del o_pro["_id"]
                s_sql = _get_select_sql(s_class, s_pro)
                o_sql = _get_select_sql(o_class, o_pro)
                __command_create(relationship, s_sql,o_sql, p_pro)

    def import2orient(self):
        # self.create_ontology()
        self.create_index()
        # self.create_entity()
        self.create_relationship()

from collections import defaultdict

class NoDe(object):
    """
    node
    """
    def __init__(self, str='', is_root=False):
        self._next_p = {}
        self.fail = None
        self.is_root = is_root
        self.str = str
        self.parent = None

    def __iter__(self):
        return iter(self._next_p.keys())

    def __getitem__(self, item):
        return self._next_p[item]

    def __setitem__(self, key, value):
        _u = self._next_p.setdefault(key, value)
        _u.parent = self

    def __repr__(self):
        return "<NoDe object '%s' at %s>" % \
               (self.str, object.__repr__(self)[1:-1].split('at')[-1])

    def __str__(self):
        return self.__repr__()


class AhoCorasick(object):
    """
    Ac object
    """
    def __init__(self, words):
        self.words_set = set(words)
        self.words = list(self.words_set)
        self.words.sort(key=lambda x: len(x))
        self._root = NoDe(is_root=True)
        self._node_meta = defaultdict(set)
        self._node_all = [(0, self._root)]
        _a = {}
        for word in self.words:
            for w in word:
                _a.setdefault(w, set())
                _a[w].add(word)

        def node_append(keyword):
            assert len(keyword) > 0
            _ = self._root
            for _i, k in enumerate(keyword):
                node = NoDe(k)
                if k not in _:
                    _[k] = node
                    self._node_all.append((_i+1, _[k]))
                if _i >= 1:
                    for _j in _a[k]:
                        if keyword[:_i+1].endswith(_j):
                            self._node_meta[id(_[k])].add((_j, len(_j)))
                _ = _[k]
            # else:
            #     if _ != self._root:
            #         self._node_meta[id(_)].add((keyword, len(keyword)))

        for word in self.words:
            node_append(word)
        self._node_all.sort(key=lambda x: x[0])
        self._make()

    def _make(self):
        """
        build ac tree
        :return:
        """
        for _level, node in self._node_all:
            if node == self._root or _level <= 1:
                node.fail = self._root
            else:
                _node = node.parent.fail
                while True:
                    if node.str in _node:
                        node.fail = _node[node.str]
                        break
                    else:
                        if _node == self._root:
                            node.fail = self._root
                            break
                        else:
                            _node = _node.fail

    def search(self, content, with_index=False):
        result = set()
        node = self._root
        index = 0
        for i in content:
            while 1:
                if i not in node:
                    if node == self._root:
                        break
                    else:
                        node = node.fail
                else:
                    for keyword, keyword_len in self._node_meta.get(id(node[i]), set()):
                        if not with_index:
                            result.add(keyword)
                        else:
                            result.add((keyword, (index - keyword_len + 1, index + 1)))
                    node = node[i]
                    break
            index += 1
        return result

    # 用于解释模型
    def search_filter_by_indx(self, content, with_index=False, v_index=0):
        result = set()
        node = self._root
        index = 0
        for i in content:
            while 1:
                if i not in node:
                    if node == self._root:
                        break
                    else:
                        node = node.fail
                else:
                    for keyword, keyword_len in self._node_meta.get(id(node[i]), set()):
                        if not with_index:
                            result.add(keyword)
                        else:
                            result.add((keyword, (index - keyword_len + 1, index + 1)))
                    node = node[i]
                    break
            index += 1
            if index >= v_index:
                break
        return result


# # 切割字符串，生成器形式
# def split_str(input_str: str, split_char: Set = None, not_add_char: Set = None):
#     if split_char is None:
#         split_char = {"。", "？", "！", "!", "?", "\r", "\n", " "}
#     if not_add_char is None:
#         not_add_char = {"\r", "\n", " "}
#     start = 0
#     for i, i_char in enumerate(input_str):
#         if i_char not in split_char:
#             continue
#         if i > start:
#             if i_char in not_add_char:
#                 yield input_str[start:i]
#             else:
#                 yield input_str[start:i+1]
#         start = i+1
#     if start < len(input_str):
#         yield input_str[start:]


class RuleValueTree(object):

    def __init__(self, value=None):
        self.value = value
        self.next = dict()
        self.is_leaf = 0

rule_list = [
            {
                "p": "是",
                "sub_pos_full": "n",
                "sub_dep_full": "top",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "ns",
                "sub_dep_full": "nsubj",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "n-n",
                "sub_dep_full": "nn-top",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "n",
                "sub_dep_full": "nsubj",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "nx",
                "sub_dep_full": "top",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "n",
                "sub_dep_full": "dep",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 2
            },{
                "p": "是",
                "sub_pos_full": "ns",
                "sub_dep_full": "nsubj",
                "pre_dep": "conj",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 4
            },{
                "p": "是",
                "sub_pos_full": "t-t",
                "sub_dep_full": "nn-top",
                "pre_dep": "conj",
                "obj_pos": "q",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "ns-n",
                "sub_dep_full": "nn-top",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "ns-n",
                "sub_dep_full": "nn-nsubj",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            },{
                "p": "是",
                "sub_pos_full": "n",
                "sub_dep_full": "dep",
                "pre_dep": "root",
                "obj_pos": "n",
                "obj_dep": "attr",
                "pre_sub_dis": 1
            },{
                "p": "是",
                "sub_pos_full": "n",
                "sub_dep_full": "top",
                "pre_dep": "root",
                "obj_pos": "vn",
                "obj_dep": "attr",
                "pre_sub_dis": 0
            }]


class EntityDescribeExtractByRoleAnalysis(object):
    """
        基于hanlp 进行实体和实体描述信息抽取
    """

    def __init__(self, concept_rule_path=None):
        import hanlp
        # model_path = "/root/model/Anysharedocumentmodel/hanlp/mtl/close_tok_pos_ner_srl_dep_sdp_con_electra_small_20210111_124159"
        # hanlp.pretrained.mtl.CLOSE_TOK_POS_NER_SRL_DEP_SDP_CON_ELECTRA_SMALL_ZH
        self.HanLP = hanlp.load(hanlp.pretrained.mtl.CLOSE_TOK_POS_NER_SRL_DEP_SDP_CON_ELECTRA_SMALL_ZH)
        # self.filter_p = {"是"}
        self.role_list = ["ARG0", "PRED", "ARG1"]

        input_rule_lists = []
        self.rule_root = RuleValueTree()
        if concept_rule_path and os.path.exists(concept_rule_path):
            concept_rule_df = pd.read_csv(concept_rule_path)
            for _, item_row in concept_rule_df.iterrows():
                input_rule_lists.append([item_row["p"],
                                         item_row["sub_pos_full"],
                                         item_row["sub_dep_full"],
                                         item_row["pre_dep"],
                                         item_row["obj_last_pos"],
                                         item_row["obj_last_dep"],
                                         item_row["pre_sub_dis"]])
        else:
            input_rule_lists = [[rule["p"],
                                rule["sub_pos_full"],
                                rule["sub_dep_full"],
                                rule["pre_dep"],
                                rule["obj_pos"],
                                rule["obj_dep"],
                                rule["pre_sub_dis"]] for rule in rule_list]
        logger.info("rule num {}".format(len(input_rule_lists)))
        self.build_rule_tree(input_rule_lists)

    def add_rule(self, input_rule_list):
        cur = self.rule_root
        for rule_value in input_rule_list:
            if rule_value not in cur.next:
                cur.next[rule_value] = RuleValueTree(rule_value)
            cur = cur.next[rule_value]
        cur.is_leaf = 1

    def build_rule_tree(self, input_rule_lists):
        for input_rule_list in input_rule_lists:
            self.add_rule(input_rule_list)

    def single_sentence(self, input_sentence):
        document = self.HanLP([input_sentence])
        doc_srl = document["srl"][0]
        doc_pos = document["pos/pku"][0]
        doc_word = document["tok/fine"][0]
        doc_dep = document["dep"][0]

        spo_list = self.single_document(doc_srl, doc_pos, doc_word, doc_dep)
        return spo_list

    def rule(self, spo, input_document_pos , input_document_dep, input_document_word):
        # rule_use = rule_list[0]
        sentence_dep = input_document_dep
        sentence_pos = input_document_pos
        sub = spo["s"]
        obj = spo["o"]
        pre = spo["p"]
        if pre[0] != "是":
            return True

        if spo["s"][0][-1] == "的":
            return True
        if spo["s"][0][:2] in ["他的", "你的", "我的", "它的", "她的"]:
            return True
        if spo["s"][0][:2] in ["这里", "主要", "常见", "这些"]:
            return True

        s_filter1 = ["下图", "图片", "照片", "图", "右图", "下图示意的", "左图", "表格", "右下图", "图示的", "下图反映的", "图甲",
                     "丙图", "甲图", "已图", "列表", "下表", "乙图"] + ["{}图".format(x) for x in "ABCDEFGHIJKLMNOPQRSTUVWXYZ"]
        s_filter2 = list("ABCDEFGHIJKLMNOPQRSTUVWXYZ") + list("abcdefghijklmnopqrstuvwxyz") + list("甲乙丙丁戊戌")
        s_filter3 = ["正确答案", "目的", "主要原因", "答案", "原因", "理由", "人为原因", "主要因素", "主要成分", "核心", "成因",
                     "内因", "判断理由", "实质", "目标", "主要问题", "决定性因素", "本质", "主要农作物", "特色农作物"]
        s_filter4 = ["我国", "指的", "这", "白色部分", "黄色部分", "他们", "它", "我", "他", "第一部分", "第二部分", "蓝色代表的",
                     "什么", "六部分", "意思", "本文", "下面", "你", "以下", "这里", "它们", "常见地区", "地区", "二者", "以上",
                     "某地", "之一", "她", "他们", "她们", "它们", "大家", "大伙", "我们", "下", "那", "本题", "其", "他的意思",
                     "前者", "这些", "本身", "主要业务", "如下", "右边", "左边", "选项"] + \
                    ["甲处", "乙处", "丙处", "丁处"] + ["{}处".format(x) for x in "ABCDEFGHIJKLMNOPQRSTUVWXYZ"] + \
                    ["其{}".format(x) for x in "一二三四五六七八九十"] + \
                    ["第{}方面".format(x) for x in "一二三四五六七八九十"] + \
                    ["{}处".format(x) for x in "abcdefghijklmnopqrstuvwxyz"] + \
                    ["{}区域".format(x) for x in "ABCDEFGHIJKLMNOPQRSTUVWXYZ"] + \
                    ["{}国".format(x) for x in "甲乙丙丁戊戌"] + \
                    ["{}国".format(x) for x in "ABCDEFGHIJKLMNOPQRSTUVWXYZ"]
        s_filter5 = ["现在", "此时", "后期", "六七十年代", "今日", "当年"]
        s_filter6 = ["关键", "措施", "一般规律", "特点", "产品", "反映的", "重点", "借鉴的", "反映的", "不同点", "好处",
                     "下列问题", "劣势", "走过的", "缺点", "结论", "图表", "办法", "结果", "问题", "故事讲述的", "疑点", "起因",
                     "特色", "说明的", "人为因素", "优点", "优势", "难点", "方法", "处", "题", "年", "主要措施", "个", "地面",
                     "国", "市", "区", "流", "事实", "作用", "背后", "您", "其中", "哪些"]

        filters = s_filter1 + s_filter2 + s_filter3 + s_filter4 + s_filter5 + s_filter6

        if spo["s"][0] in filters:
            return True
        if spo["o"][0][-2:] in ["如此"]:
            return True

        cur = self.rule_root
        predicate = pre[0]
        if predicate not in cur.next:
            return True
        cur = cur.next[predicate]
        sub_pos_full = "-".join(sentence_pos[sub[1]:sub[2]])
        if sub_pos_full not in cur.next:
            return True
        cur = cur.next[sub_pos_full]
        sub_dep = [item[1] for item in sentence_dep[sub[1]:sub[2]]]
        sub_dep_full = "-".join(sub_dep)
        if sub_dep_full not in cur.next:
            return True
        cur = cur.next[sub_dep_full]
        pre_dep = sentence_dep[pre[1]][1]
        if pre_dep not in cur.next:
            return True
        cur = cur.next[pre_dep]
        obj_pos = sentence_pos[obj[2] - 1]
        if obj_pos not in cur.next:
            return True
        cur = cur.next[obj_pos]
        obj_dep = sentence_dep[obj[2] - 1][1]
        if obj_dep not in cur.next:
            return True
        cur = cur.next[obj_dep]
        pre_sub_dis = pre[1] - sub[2]
        if pre_sub_dis not in cur.next:
            return True

        return False

    def single_document(self, input_document_srl, input_document_pos, input_document_word, input_document_dep=None):
        spo_list = []

        for role_part in input_document_srl:
            spo = {}
            if len(role_part) < 3:
                continue

            for role_mention, role_key, start, end in role_part:
                if role_key not in self.role_list:
                    continue
                if role_key == "PRED":
                    spo["p"] = (role_mention, start, end)
                elif role_key == "ARG0":
                    spo["s"] = (role_mention, start, end)
                elif role_key == "ARG1":
                    spo["o"] = (role_mention, start, end)

            if "p" not in spo:
                continue
            if "o" not in spo:
                continue
            if "s" not in spo:
                continue
            if spo["p"][1] < spo["s"][2]:
                continue
            if spo["p"][2] > spo["o"][1]:
                continue

            # if (spo["o"][2]-spo["s"][1])*1.0/len(input_document_word) < 0.6:
            #     # print(spo["o"], len(input_document_word))
            #     continue

            if self.rule(spo, input_document_pos, input_document_dep, input_document_word):
                continue

            spo_list.append(spo)
            break
        return spo_list

    def extract_info(self, input_data):
        if isinstance(input_data, str):
            input_sentence_list = re.split("([。？！\s+/\n])", input_data)
        elif isinstance(input_data, list):
            input_sentence_list = input_data
        else:
            raise TypeError

        # def simple_filter(i_sentence):
        #     if len(i_sentence) < 10:
        #         return False
        #     if len(i_sentence) > 100:
        #         return False
        #     if i_sentence[-1] in ["?", "？"]:
        #         return False
        #     if "是我" in i_sentence:
        #         return False
        #     if "是你" in i_sentence:
        #         return False
        #     if i_sentence[0] == "但":
        #         return False
        #     if not re.fullmatch("^[\u4e00-\u9fa5_a-zA-Z]{1,15}是.+$", i_sentence):
        #         return False
        #     if re.fullmatch("^.+吗$", i_sentence):
        #         return False
        #     return True

        # input_sentence_list = [sentence.strip() for sentence in input_sentence_list if simple_filter(sentence.strip())]
        output_documents = self.HanLP(input_sentence_list, batch_size=64, tasks=["srl", "tok/fine", "pos/pku", "dep"])
        entity_describe_res = []
        for i, document in enumerate(output_documents.get("srl", list())):
            sentence_words = output_documents["tok/fine"][i]
            out_spo_list = self.single_document(document,
                                                output_documents["pos/pku"][i],
                                                sentence_words,
                                                output_documents["dep"][i])
            for spo in out_spo_list:
                entity_describe_res.append(({"sentence": input_sentence_list[i],
                                            "entity": spo["s"][0],
                                            "describe": "".join(sentence_words[spo["s"][1]:spo["o"][2]])},
                                            i))

        return entity_describe_res


logging.basicConfig(level = logging.INFO,format = '%(asctime)s - %(name)s - %(levelname)s - %(lineno)d - %(message)s')
logger = logging.getLogger(__name__)


DELETE_STATUS = 0
ADD_STATUS = 1
UPDATE_STATUS = 2
RENAME_STATUS = 3
AS_ES_CONTENT_FAIL_STATUS = 4

DIR_TYPE = 1
FILE_TYPE = 0
data_lock = threading.Lock()

import torch
# import multiprocessing
# cpu_core_count = multiprocessing.cpu_count()
# if cpu_core_count > 1:
#     cpu_core_count -= 1
#
# # 限制cpu 核数
# torch.set_num_threads(1)
import multiprocessing
from multiprocessing import Pool, Process, Manager


# 概念解释模型批处理：包括找到符合条件的句子和插入mongodb
def process_tim(model, input_sentence_list, input_infos, ds_id, db, collection_name_map, entity_property_dict, relation_property_dict):
    labels = []
    label2documents = []
    desc2label = []
    desc2document = []
    desc = []
    res = model.extract_info(input_sentence_list)
    # print("res num {}".format(len(res)))
    for rs, iv in res:
        if "desc2label" in relation_property_dict:
            desc2label.append({"s": rs["describe"], "p": "desc2label", "o": rs["entity"],
                               "p_pro": {"name": "desc2label"}, "s_pro": {"name": rs["describe"]},
                               "o_pro": {"name": rs["entity"], "adlabel_kcid": rs["entity"]}})
        if "desc2document" in relation_property_dict:
            desc2document.append({"s": rs["describe"], "p": "desc2document", "o": input_infos[iv]["name"],
                                  "s_pro": {"name": rs["describe"], "ds_id": ds_id},
                                  "p_pro": {"name": "desc2document"},
                                  "o_pro": {"gns": input_infos[iv]["gns"]}})
        if "desc" in entity_property_dict:
            desc.append({"name": rs["describe"], "ds_id": ds_id})
        if "label" in entity_property_dict:
            labels.append({
                "name": rs["entity"],
                "adlabel_kcid": rs["entity"],
                "kc_topic_tags": "",
                "confidence": "",
                "type_sa": "true",
                "ds_id": ds_id
            })
        if "label2documents" in relation_property_dict:
            label2documents.append({
                "s": rs["entity"],
                "p": "label2document",
                "o": input_infos[iv]["name"],
                "s_pro": {
                    "name": rs["entity"],
                    "adlabel_kcid": rs["entity"],
                    "ds_id": ds_id,
                    "type_sa": "true"
                },
                "p_pro": {
                    "name": "label2document"
                },
                "o_pro": {
                    "gns": input_infos[iv]["gns"]
                }

            })
    if desc2label:
        db[collection_name_map["desc2label"]].insert_many(desc2label)
    if desc2document:
        db[collection_name_map["desc2document"]].insert_many(desc2document)
    if desc:
        db[collection_name_map["desc"]].insert_many(desc)
    if labels:
        db[collection_name_map["label"]].insert_many(labels)
    if label2documents:
        db[collection_name_map["label2document"]].insert_many(label2documents)


# 单进程概念解释处理模式
def single_process_model(hanlp_model_path, concept_rule_path, ds_id, graph_name, batch_num, q, dt,
                          collection_name_map,
                         entity_property_dict={}, relation_property_dict={}):
    db = MongoClientAD().connect_mongo()
    os.environ['HANLP_HOME'] = hanlp_model_path
    model = EntityDescribeExtractByRoleAnalysis(concept_rule_path=concept_rule_path)

    print("sub process start")
    while True:
        if not q.empty():
            value = q.get(True)
            if value["state"] == 1:
                break
            _id = value["_id"]
            if _id is None:
                items = db[collection_name_map["temp_concept"]].find().limit(batch_num).sort("_id", 1)
            else:
                items = db[collection_name_map["temp_concept"]].find({"_id": {"$gt": _id}}).limit(batch_num).sort("_id", 1)
            input_sentence_list = []
            input_infos = []
            for item in items:
                input_sentence_list.append(item["sentence"])
                input_infos.append({"gns": item["gns"],
                                    "name": item["name"]})
            state = 0
            try:
                process_tim(model, input_sentence_list, input_infos, ds_id, db,
                            collection_name_map,
                            entity_property_dict,
                            relation_property_dict)
            except Exception as e:
                logger.error(e)
                state = 1

            dt["finish"] += len(input_sentence_list)
            if state:
                logger.error("anyshare concept {}/{} fail".format(dt["finish"], dt["count"]))
            else:
                logger.info("anyshare concept {}/{} finish".format(dt["finish"], dt["count"]))

            del input_sentence_list
            gc.collect()
        else:
            time.sleep(3)
    del model
    del db
    gc.collect()

    print("sub process end")


class AnyshareDocumentModel(object):
    """
        文档结构模型 精简版
        用法

        asdm = AnyshareDocumentModel(entity_property_dict, relationship_property_dict, gnsid, conn, config_path, as7json, graph_name, None, None, None, None)
        asdm.start()
    """
    def __init__(self, entity_property_dict, relationship_property_dict, conn_db, as7json,is_increment=None, config_path=None, **kwargs):
        # self.conn = conn
        graph_collection_prefix = kwargs.get("graph_collection_prefix", "")

        self.gnsid = as7json["docid"]
        self.graph_name = graph_collection_prefix
        self.mongo_gns_info_db = conn_db
        self.as7json = as7json
        self.ds_auth = as7json["ds_auth"]
        self.ds_address = as7json["ds_address"]
        self.ds_port = as7json["ds_port"]
        self.ds_name = as7json["dsname"]
        self.ds_id = as7json["id"]
        self.version = None
        self.entity_property_dict = entity_property_dict
        self.relationship_property_dict = relationship_property_dict
        self.sess = requests.Session()
        self.is_increment = is_increment
        self.config_path = config_path
        self.debug = False

        # 设置hanlp环境变量，要放在解释模型之前，不然不会生效
        self.hanlp_model_path = self.get_config()["concept"]["hanlp_model_path"]
        os.environ['HANLP_HOME'] = self.hanlp_model_path
        self.concept_rule_path = self.get_config()["concept"]["concept_description_rule_path"]
        self.ede_model = None

        # 概念解释专用
        self.data_cache = []
        self.graph_collection_prefix = graph_collection_prefix

        self.collection_name_map = {
            "gnsInfo": "{}_gnsInfo".format(self.graph_collection_prefix),
            "gnsInfoChange": "{}_gnsInfoChange".format(self.graph_collection_prefix),
            "chapter": "{}_chapter".format(self.graph_collection_prefix),
            "chapter2text": "{}_chapter2text".format(self.graph_collection_prefix),
            "desc": "{}_desc".format(self.graph_collection_prefix),
            "desc2document": "{}_desc2document".format(self.graph_collection_prefix),
            "desc2label": "{}_desc2label".format(self.graph_collection_prefix),
            "document": "{}_document".format(self.graph_collection_prefix),
            "document2chapter": "{}_document2chapter".format(self.graph_collection_prefix),
            "document2text": "{}_document2text".format(self.graph_collection_prefix),
            "folder": "{}_folder".format(self.graph_collection_prefix),
            "folder2document": "{}_folder2document".format(self.graph_collection_prefix),
            "folder2folder": "{}_folder2folder".format(self.graph_collection_prefix),
            "label": "{}_label".format(self.graph_collection_prefix),
            "label2document": "{}_label2document".format(self.graph_collection_prefix),
            "label2label": "{}_label2label".format(self.graph_collection_prefix),
            "subject": "{}_subject".format(self.graph_collection_prefix),
            "subject2document": "{}_subject2document".format(self.graph_collection_prefix),
            "text2text": "{}_text2text".format(self.graph_collection_prefix),
            "text": "{}_text".format(self.graph_collection_prefix),
            "temp_concept": "{}_temp_concept".format(self.graph_collection_prefix),
            "empty_data": "{}_empty_data".format(self.graph_collection_prefix),
        }

    def get_config(self):
        config = ConfigParser()
        config.read(self.config_path, "utf-8")
        return config

    # 检查索引，如果没有就创建
    # def check_and_create(self):
    #     data_index = self.mongo_gns_info_db["label2document"].index_information()
    #     status = 1
    #     for k, _ in data_index.items():
    #         if k.startswith("o_pro.gns"):
    #             status = 0
    #             break
    #     if status:
    #         self.mongo_gns_info_db["label2document"].create_index([('o_pro.gns',1)])
    #
    #     data_index = self.mongo_gns_info_db["label"].index_information()
    #     status = 1
    #     for k, _ in data_index.items():
    #         if k.startswith("name"):
    #             status = 0
    #             break
    #     if status:
    #         self.mongo_gns_info_db["label"].create_index([('name', 1)])

    # 获取gns info 信息
    def get_all_gns(self):
        # data_collection = "gnsInfo"
        # if self.is_increment:
        data_collection = self.collection_name_map["gnsInfoChange"]
        g_info_list = self.mongo_gns_info_db[data_collection].find({}, {"_id": 0})
        for g_info in g_info_list:
            yield g_info

        # g_info_list = self.mongo_gns_info_db["empty_data"].find({}, {"_id": 0})
        # for g_info in g_info_list:
        #     yield g_info

    # 获取指定文件的tag
    def get_tags_by_gns(self, gns, ds_address, ds_port, sess: requests.Session = None):
        print('开始获取AS token', __file__, 'get_tags_by_gns')
        ret_code, acc_token = asToken.get_token(self.ds_auth)
        item_id = gns.split("/")[-1]
        url = "https://{0}:{1}/api/metadata/v1/item/{2}/tag".format(ds_address, ds_port, item_id)
        headers = {
            'Content-Type': 'application/json',
            'Authorization': 'Bearer ' + acc_token,
        }
        res = sess.get(url, headers=headers, timeout=10, verify=False)
        if res.status_code != 200:
            logger.error(res.content)
            return {"tag": []}

        return res.json()["tag"]

    # 获取文件内容， 基于as 文件下载接口
    def get_as_file_content(self, gns, ds_address, ds_port, sess: requests.Session = None):
        print('开始获取AS token', __file__, 'get_as_file_content')
        ret_code, acc_token = asToken.get_token(self.ds_auth)
        as_url = "{}:{}/api/efast/v1/file/osdownload".format(ds_address, ds_port)
        headers = {
            'Content-Type': 'application/json',
            'Authorization': 'Bearer ' + acc_token,
        }
        response = sess.post(url=as_url, headers=headers, json={"docid": gns}, timeout=10, verify=False)
        if response.status_code != 200:
            logger.error(response.content)
            return ""
        authrequest = response.json()["authrequest"]
        one_headers = {}
        for i in authrequest[2:]:
            one_headers[i[:i.index(":")].strip()] = i[i.index(":") + 1:].strip()
        response = sess.request(authrequest[0], authrequest[1], headers=one_headers, verify=False)
        if response.status_code != 200:
            logger.error(response.content)
            return ""
        return response.content

    # 解析文本内容
    def parse_file_content(self, input_file_content):
        head_property_list = [item[0] for item in self.entity_property_dict["chapter"]]
        text_property_list = [item[0] for item in self.entity_property_dict["text"]]

        content = input_file_content["content"]
        text_len = 1000
        if content:
            # 测试用的
            # print(type(content))
            # self.mongo_gns_info_db["content"].insert_one({
            #     "gns": input_file_content["gns"],
            #     "name": input_file_content["name"],
            #     "content": content
            # })
            content = content.strip()
            chapter = re.split(r"[\s\r\n]", content)[0].strip()
            content = re.sub(r"\s+", " ", content)
            # content = content.replace(" ", "")

            chapter_info = dict()
            # if input_file_content["file_type"] not in [".txt"]:
            #     self.mongo_gns_info_db["content"].insert_one({"gns": input_file_content["gns"], "content": content})
            if len(chapter) > 20:
                chapter = None
            else:
                chapter_info = {
                    head_property_list[0]: chapter,
                    head_property_list[1]: input_file_content["path"],
                    head_property_list[2]: 1,
                    "ds_id": self.ds_id
                }
                if "chapter" in self.entity_property_dict:
                    self.mongo_gns_info_db[self.collection_name_map["chapter"]].insert_one(chapter_info)
                if "document2chapter" in self.relationship_property_dict:
                    doc2chapter_info = {"s": input_file_content["name"],
                                        "p": "document2chapter",
                                        "o": chapter,
                                        "s_pro": {"gns": input_file_content["gns"]},
                                        "p_pro": {"name": "document2chapter"},
                                        "o_pro": chapter_info,
                                        "ds_id": self.ds_id}
                    self.mongo_gns_info_db[self.collection_name_map["document2chapter"]].insert_one(doc2chapter_info)

            text_list = []
            text2text = []
            chapter2text = []
            doc2text = []
            last = None

            text_span_id = 0
            while text_span_id < len(content):
                while text_span_id < len(content) and content[text_span_id] in {"。", "!", ";", ".", "?", "？", "！", ",", "，", "；"}:
                    text_span_id += 1

                last_text_ind = text_span_id+text_len
                while len(content)> last_text_ind \
                        and last_text_ind > text_span_id \
                        and content[last_text_ind-1] not in {",", "\n", "。", "!","\r", ".", "?", "？", "！", "，", ";"}:
                    last_text_ind -= 1
                if last_text_ind == text_span_id:
                    last_text_ind = text_span_id+text_len

                text_span = content[text_span_id:last_text_ind].strip()
                text_span_id = last_text_ind
                if text_span == "":
                    continue
                text_list.append({text_property_list[0]: text_span, "ds_id": self.ds_id})
                doc2text.append({"s": input_file_content["name"], "p": "document2text", "o": text_span,
                                 "s_pro": {"gns": input_file_content["gns"]},
                                 "p_pro": {"name": "document2text"},
                                 "o_pro": {"name": text_span, "ds_id": self.ds_id}})
                if last:
                    text2text.append({"s": last, "p": "text2text", "o": text_span,
                                      "s_pro": {"name": last},
                                      "p_pro": {"name": "text2text"},
                                      "o_pro": {"name": text_span},
                                      "ds_id": self.ds_id})

                if chapter:
                    chapter2text.append({"s": chapter, "p": "chapter2text", "o": text_span,
                                         "s_pro": chapter_info,
                                         "p_pro": {"name": "chapter2text"},
                                         "o_pro": {"name": text_span, "ds_id": self.ds_id},
                                         "ds_id": self.ds_id})
                last = text_span

            if text_list and "text" in self.entity_property_dict:
                self.mongo_gns_info_db[self.collection_name_map["text"]].insert_many(text_list)
            if text2text and "text2text" in self.relationship_property_dict:
                self.mongo_gns_info_db[self.collection_name_map["text2text"]].insert_many(text2text)
            if doc2text and "document2text" in self.relationship_property_dict:
                self.mongo_gns_info_db[self.collection_name_map["document2text"]].insert_many(doc2text)
            if chapter2text and "chapter2text" in self.relationship_property_dict:
                self.mongo_gns_info_db[self.collection_name_map["chapter2text"]].insert_many(chapter2text)
            if "desc" in self.entity_property_dict or "label" in self.entity_property_dict:
                self.add_entity_describe(content, input_file_content)

    # 增加实体解释
    def add_entity_describe(self, input_content, input_file_info, newword_label_set=None):
        input_sentence_list = []
        for sentence in split_str(input_content):

            input_sentence_list.append({"sentence": sentence,
                                        "gns": input_file_info["gns"],
                                        "name": input_file_info["name"]})
        #
        if input_sentence_list:
            self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].insert_many(input_sentence_list)

    def insert_mongo_task(self, collection_name, data):
        # print(collection_name, len(data))
        self.mongo_gns_info_db[self.collection_name_map[collection_name]].insert_many(data)

    def update_mongo_task(self, collection_name, label_updates):
        for item in label_updates:
            self.mongo_gns_info_db[self.collection_name_map[collection_name]].update_many({"name": item},
                                                        {"$set": {"type_sa": "true"}})

    # 多线程批量处理概念解释句子数据
    def batch_entity_describe_v2(self):

        t_pool = ThreadPoolExecutor(max_workers=100)
        f_list = []
        first = True
        last_id = None
        data_all = self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].count()
        data_collection_num = 0
        batch_num = 64
        # 使用limit $gt _id 这种遍历方式
        while True:
            # print(last_id)
            input_sentence_list = []
            input_infos = []
            labels = []
            label2documents = []
            desc2label = []
            desc2document = []
            desc = []
            # start_time = time.time()
            if first:
                items = self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].find().limit(64).sort("_id", 1)
                first = False
            else:
                items = self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].find({"_id": {"$gt": last_id}}).limit(64).sort("_id", 1)

            for item in items:
                input_sentence_list.append(item["sentence"])
                input_infos.append(item)

            res = self.ede_model.extract_info(input_sentence_list)
            for rs, iv in res:

                desc2label.append({"s": rs["describe"], "p": "desc2label", "o": rs["entity"],
                                   "p_pro": {"name": "desc2label"}, "s_pro": {"name": rs["describe"]},
                                   "o_pro": {"name": rs["entity"], "adlabel_kcid": rs["entity"]}})

                desc2document.append({"s": rs["describe"], "p": "desc2document", "o": input_infos[iv]["name"],
                                      "s_pro": {"name": rs["describe"], "ds_id": self.ds_id},
                                      "p_pro": {"name": "desc2document"},
                                      "o_pro": {"gns": input_infos[iv]["gns"]}})
                desc.append({"name": rs["describe"], "ds_id": self.ds_id})

                labels.append({
                    "name": rs["entity"],
                    "adlabel_kcid": rs["entity"],
                    "kc_topic_tags": "",
                    "confidence": "",
                    "type_sa": "true",
                    "ds_id": self.ds_id
                })
                label2documents.append({
                    "s": rs["entity"],
                    "p": "label2document",
                    "o": input_infos[iv]["name"],
                    "s_pro": {
                        "name": rs["entity"],
                        "adlabel_kcid": rs["entity"],
                        "ds_id": self.ds_id,
                        "type_sa": "true"
                    },
                    "p_pro": {
                        "name": "label2document"
                    },
                    "o_pro": {
                        "gns": input_infos[iv]["gns"]
                    }

                    })
            if desc2label and "desc2label" in self.relationship_property_dict:
                th = t_pool.submit(self.insert_mongo_task, "desc2label", desc2label)
                f_list.append(th)
            if desc2document and "desc2document" in self.relationship_property_dict:
                th = t_pool.submit(self.insert_mongo_task, "desc2document", desc2document)
                f_list.append(th)
            if desc and "desc" in self.entity_property_dict:
                th = t_pool.submit(self.insert_mongo_task, "desc", desc)
                f_list.append(th)
            if labels and "label" in self.entity_property_dict:
                th = t_pool.submit(self.insert_mongo_task, "label", labels)
                f_list.append(th)
            if label2documents and "label2document" in self.relationship_property_dict:
                th = t_pool.submit(self.insert_mongo_task, "label2document", label2documents)
                f_list.append(th)
            data_collection_num += len(input_sentence_list)
            logger.info("desc extract load {}/{}".format(data_collection_num, data_all))
            if len(input_sentence_list) < batch_num:
                break
            last_id = input_infos[-1]["_id"]
            # print("batch cost {}".format(time.time()-start_time))

        wait(f_list)
        self.ede_model = None
        gc.collect()

        if not self.debug:
            self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].drop()

    # 多进程方式
    def batch_entity_describe_v4(self):
        import psutil
        process_num = multiprocessing.cpu_count()
        if process_num > 1:
            process_num -= 1

        # 根据内存分配核心数,每个进程大概占1.4G
        mem = psutil.virtual_memory()
        total_mem = float(mem.total) / 1024 / 1024 / 1024
        available_mem = float(mem.available) / 1024 / 1024 / 1024

        logger.info("anyshare concept desc task: cpu num:{0}".format(multiprocessing.cpu_count()))
        logger.info("anyshare concept desc task: memory total:{0}".format(total_mem))
        logger.info("anyshare concept desc task: memory available:{0}".format(available_mem))

        process_num = min(int(available_mem//1.4), process_num)

        # 如果cpu 核心数量比较少，就用多线程的方式
        if process_num < 4:
            self.ede_model = EntityDescribeExtractByRoleAnalysis(concept_rule_path=self.concept_rule_path)
            self.batch_entity_describe_v2()
            return
        torch.set_num_threads(1)

        # data_all = self.mongo_gns_info_db["temp_concept"].count()
        data_all = self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].count()
        if data_all == 0:
            return
        manager = Manager()
        dt = manager.dict()
        dt["count"] = data_all
        dt["finish"] = 0
        batch_size = 256
        batch_num = int(data_all//batch_size) if data_all % batch_size == 0 else int(data_all//batch_size)+1
        if batch_num < process_num:
            process_num = batch_num

        logger.info("anyshare concept desc task: process num:{0}".format(process_num))

        queue = multiprocessing.Queue(10000000)
        p_list = [Process(target=single_process_model, args=(self.hanlp_model_path, self.concept_rule_path, self.ds_id,
                                                             self.graph_name, batch_size, queue, dt,
                                                             self.collection_name_map,
                                                             self.entity_property_dict,
                                                             self.relationship_property_dict))
                  for _ in range(process_num)]
        [task.start() for task in p_list]
        first = True
        last_id = None

        # 使用limit $gt _id 这种遍历方式
        while True:
            # print(last_id)
            # input_sentence_list = []
            # input_infos = []

            if first:
                items = self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].find().limit(batch_size).sort("_id", 1)
                queue.put({
                    "_id": None,
                    "state": 0
                })
                first = False
            else:
                queue.put({
                    "_id": last_id,
                    "state": 0
                })
                items = self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].find({"_id": {"$gt": last_id}}).limit(batch_size).sort(
                    "_id", 1)
            data_list = [item["_id"] for item in items]
            if len(data_list) < batch_size:
                break
            last_id = data_list[-1]

        [queue.put({"state": 1}) for _ in range(process_num)]

        [task.join() for task in p_list]

        if not self.debug:
            self.mongo_gns_info_db[self.collection_name_map["temp_concept"]].drop()

        del queue
        del manager
        gc.collect()

    # 对目录信息处理
    def task_folder(self, gns_info):
        folder_property_list = [item[0] for item in self.entity_property_dict["folder"]]
        self.mongo_gns_info_db[self.collection_name_map["folder"]].insert_one({folder_property_list[0]: gns_info["name"],
                                                     folder_property_list[1]: gns_info["path"],
                                                     folder_property_list[2]: gns_info["gns"],
                                                     folder_property_list[3]: gns_info["create_time"],
                                                     "rev": gns_info["rev"],
                                                     "ds_id": self.ds_id})
        if "folder2document" in self.relationship_property_dict:
            for child_file in gns_info["child_file_info"]:
                insert_info = {
                    "s": gns_info["name"],
                    "p": "folder2document",
                    "o": child_file["name"],
                    "s_pro": {"gns": gns_info["gns"]},
                    "o_pro": {"gns": child_file["gns"]},
                    "p_pro": {"name": "folder2document", "domain": child_file["name"].split(".")[-1],
                              "range": child_file["name"]},
                    "ds_id": self.ds_id
                }
                self.mongo_gns_info_db[self.collection_name_map["folder2document"]].insert_one(insert_info)
        if "folder2folder" in self.relationship_property_dict:
            for child_file in gns_info["child_folder_info"]:
                insert_info = {
                    "s": gns_info["name"],
                    "p": "folder2folder",
                    "o": child_file["name"],
                    "s_pro": {"gns": gns_info["gns"]},
                    "o_pro": {"gns": child_file["gns"]},
                    "p_pro": {"name": "folder2folder", "domain": gns_info["name"], "range": child_file["name"]},
                    "ds_id": self.ds_id
                }

                self.mongo_gns_info_db[self.collection_name_map["folder2folder"]].insert_one(insert_info)

    # 对document 信息处理
    def task_document(self, gns_info):
        # 版本比较函数， todo 临时存放点
        def version_greater_equal(version_a, version_b):
            version_a = version_a.split(".")
            version_b = version_b.split(".")

            for i, va in enumerate(version_a):
                if int(va) < int(version_b[i]):
                    return False
                elif int(va) > int(version_b[i]):
                    return True
            return True
        document_property_list = [item[0] for item in self.entity_property_dict["document"]]

        file_type = gns_info["name"].split(".")[-1]
        gns_info["ds_address"] = self.ds_address
        gns_info["ds_id"] = self.ds_id
        gns_info["file_type"] = "." + file_type.lower()

        gns_info["modified_time"] = time.strftime("%Y-%m-%d %H:%M:%S",
                                                  time.localtime(int(str(gns_info["modified"])[0:10])))
        # print("start vvvvvv")
        if file_type.lower() in ["txt", "doc", "ppt", "pdf", "docx", "pptx"]:
            if self.version == "7.0.1.7":
                retcode = 200
                source = {
                    "tags": self.get_tags_by_gns(gns_info["gns"], self.ds_address, self.ds_port, self.sess),
                    "content": self.get_as_file_content(gns_info["gns"], self.ds_address, self.ds_port, self.sess)
                }

            else:
                acc_token = None
                if version_greater_equal(self.version, "7.0.1.8"):
                    print('开始获取AS token', __file__, 'task_document')
                    ret_code, acc_token = asToken.get_token(self.ds_auth)

                # logger.info(acc_token)
                retcode, source = otl_dao.getinfofromas_v2(self.ds_id, ["content", "tags"], gns_info["gns"], self.ds_address,
                                                           self.ds_port, self.version, self.sess, acc_token)
                # print(source)
                if retcode != 200:
                    logger.error(source)
                    source = {"tags": []}
            # logger.error(source)
            if "content" not in source:
                logger.error("file [{}] fail get content".format(gns_info["path"]))
                gns_info["change_type"] = AS_ES_CONTENT_FAIL_STATUS
                if "first_faile_time" not in gns_info:
                    gns_info["first_faile_time"] = time.time()
                # 3600 * 24 * 7 一周时间，如果还是请求不到数据就放弃 ╮(╯▽╰)╭
                if time.time() - gns_info["first_faile_time"] < 604800:
                    self.mongo_gns_info_db[self.collection_name_map["empty_data"]].insert_one(gns_info)
                source["content"] = ""
                # return "(๑￫ܫ￩)"
            gns_info["label"] = source["tags"]
            gns_info["content"] = source["content"]

            # if "null" in gns_info["label"]:
            #     print(gns_info["gns"])
            #     print(gns_info["name"])
            #     print(gns_info["label"])

            labels = []
            label2documents = []
            # newword_label = {n_info["s"] for n_info in self.mongo_gns_info_db["label2document"]
            #     .find({"o_pro.gns": gns_info["gns"]}, {"s": 1})}
            # newword_label = set()
            if "label" in self.entity_property_dict:
                for label in gns_info["label"]:
                    # label_status = self.check_label(label)
                    labels.append({
                        "name": label,
                        "adlabel_kcid": label,
                        "kc_topic_tags": "",
                        "confidence": "",
                        "type_as": "true",
                        "ds_id": self.ds_id
                    })
                    if "label2document" not in self.relationship_property_dict:
                        continue
                    label2documents.append({
                        "s": label,
                        "p": "label2document",
                        "o": gns_info["name"],
                        "s_pro": {
                            "name": label,
                            "adlabel_kcid": label,
                            "ds_id": self.ds_id,
                            "type_as": "true"
                        },
                        "p_pro": {
                            "name": "label2document"
                        },
                        "o_pro": {
                            "gns": gns_info["gns"]
                        }
                    })
            if labels:
                # self.mongo_gns_info_db[self.collection_name_map["label"]].insert_many(labels)
                # for label in labels:
                    # self.label_as_set.add(label["name"])
                if "label" in self.entity_property_dict:
                    self.mongo_gns_info_db[self.collection_name_map["label"]].insert_many(labels)
            if label2documents:
                if "label2document" in self.relationship_property_dict:
                    self.mongo_gns_info_db[self.collection_name_map["label2document"]].insert_many(label2documents)
            try:
                self.parse_file_content(gns_info)
            except Exception as e:
                print(e)

        doc_info = {document_property_list[0]: gns_info["name"],  ###name
                    document_property_list[1]: gns_info["path"],  ###path
                    document_property_list[2]: gns_info["creator"],  ###creator
                    document_property_list[3]: gns_info["create_time"],  ###create_time
                    document_property_list[4]: gns_info["editor"],  ###editor
                    document_property_list[5]: gns_info["gns"],  ###gns
                    document_property_list[6]: gns_info.get("file_type", file_type),
                    document_property_list[7]: gns_info.get("modified_time", ""),
                    "rev": gns_info["rev"],
                    "ds_id": self.ds_id
                    }
        # print(doc_info)
        if "document" in self.entity_property_dict:
            self.mongo_gns_info_db[self.collection_name_map["document"]].insert_one(doc_info)

    # 多线程
    def start(self):
        logger.info("anyshare document extract start!")
        start_time = time.time()
        t_pool = ThreadPoolExecutor(max_workers=100)
        print('开始获取AS token', __file__, 'start')
        ret_token, obj_token = asToken.get_token(self.ds_auth)
        if ret_token != 200:
            logger.error(obj_token)
            raise Exception("anyshare token has something wrong!")

        # 现在不需要建立索引了
        # self.check_and_create()

        # https://www.cnblogs.com/pengyusong/p/5802929.html requests connection pool is full discarding connection问题 解决方案
        self.sess.mount('https://', HTTPAdapter(pool_connections=100, pool_maxsize=100))
        version = otl_dao.getversion(self.ds_address)
        self.version = version.split("-")[2]
        logger.info("as adress {} as port {}".format(self.ds_address, self.ds_port))
        logger.info("as version {}".format(self.version))
        f_list = []
        for gns_info in self.get_all_gns():
            if gns_info.get("change_type", None) == DELETE_STATUS:
                continue
            if "create_time" in gns_info:
                if str(gns_info["create_time"]).isdigit():
                    gns_info["create_time"] = time.strftime("%Y-%m-%d %H:%M:%S",
                                                    time.localtime(int(str(gns_info["create_time"])[0:10])))
            else:
                gns_info["create_time"] = ""
            if gns_info["type"] == DIR_TYPE:
                if "folder" not in self.entity_property_dict:
                    continue
                th = t_pool.submit(self.task_folder, gns_info)
                f_list.append(th)
            else:
                th = t_pool.submit(self.task_document, gns_info)
                f_list.append(th)
        wait(f_list)
        self.sess.close()

        fail_read_doc_num = self.mongo_gns_info_db[self.collection_name_map["empty_data"]].count()
        logger.info("document fail to read num {}".format(fail_read_doc_num))
        # self.mongo_gns_info_db["empty_data"].drop()
        # if fail_read_doc_num > 0:
        #     self.mongo_gns_info_db["empty_data2"].rename("empty_data")

        logger.info("anyshare document extract end! cost {}s".format(time.time() - start_time))
        start_time = time.time()

        # 对概念解释进行收尾
        self.batch_entity_describe_v4()
        gc.collect()
        logger.info("anyshare concept extract end! cost {}s".format(time.time()-start_time))
        # start_time = time.time()

        # 主題信息
        # self.task_pre_subject()
        # self.create_subject2document()
        # logger.info("anyshare text match task end! cost {}s".format(time.time() - start_time))
        gc.collect()


NO_CHANGE = 0
SOMETHING_CHANGE = 2
MONGODB_EMPTY = 1
LASTESDATA_EMPTY = 3


# def asd_decorator(func):
#     def wrapper(*args, **kw):
#         f_res = func(*args, **kw)
#         if "code" in f_res and f_res[""]:
#             logger.error(f_res)
#             return dict()
#         return f_res
#
#     return wrapper


def print_run_time(func):
    def wrapper(*args, **kw):
        local_time = time.time()
        res = func(*args, **kw)
        logger.info("anyshare gns info cost {:.2f}s".format(time.time() - local_time))
        logger.info("anyshare gns info construct info end")
        return res
    return wrapper


class AsShareDocumentChangeModel(object):
    """
        anyshare 文档结构增量模型
        use
        as_change_model = AsShareDocumentChangeModel(as7json, conn, graph_name)
        as_change_model.compare()
        # as_change_model.build_tree("gns://400946019E6C4238991CFA8CD73B9D69")
        # as_change_model.compare()
    """

    def __init__(self, as7json, conn_db, graph_name=None, input_otl_dao=None, graph_collection_prefix=None):
        self.as7json = as7json
        self.ds_address = as7json["ds_address"]
        self.ds_port = as7json["ds_port"]
        self.mongo_kg_db = conn_db
        if input_otl_dao is None:
            self.otl_dao = otl_dao
        else:
            self.otl_dao = input_otl_dao
        self.ds_auth = as7json["ds_auth"]
        self.graph_collection_prefix = graph_collection_prefix

        self.change_collection = self.mongo_kg_db["{}_gnsInfoChange".format(self.graph_collection_prefix)]
        self.gns_info = self.mongo_kg_db["{}_gnsInfo".format(self.graph_collection_prefix)]
        self.valid_key = ["gns", "name", "creator", "editor", "create_time", "path", "child_file", "child_folder",
                          "rev", "type", "root", "root_gns", "child_folder_info", "child_file_info", "modified"]
        self.valid_change_key = self.valid_key + ["change_type", "sub_root"]
        self.queue = queue.Queue(1000000)
        self.bk_queue = queue.Queue(1000000)

    # 获取指定目录下的文件和目录
    # @asd_decorator
    def get_dir_list(self, input_gns, fail_time=0, token_fail=False):
        if fail_time == 3:
            logger.error("get dir list fail exceed 3 times")
            return dict()
        print('开始获取AS token', __file__, 'get_dir_list')
        ret_code, token_id = asToken.get_token(self.ds_auth)
        headers = {
            'Content-Type': 'application/json',
            'Authorization': 'Bearer {}'.format(token_id)
        }

        payload = {
            "docid": input_gns
        }
        url = "{}:{}/api/efast/v1/dir/list".format(self.ds_address, self.ds_port)

        response = requests.request("POST", url, headers=headers, json=payload, timeout=100, verify=False)
        if response.status_code != 200:
            if response.status_code == 401:
                response_json = response.json()
                if response_json["code"] == 401001001:
                    print('开始获取AS token', __file__, 'get_dir_list 2')
                    asToken.get_token(self.ds_auth)

            # logger.error(response.status_code)
            logger.error("get_dir_list fail {}".format(fail_time+1))
            logger.error(response.content)
            return self.get_dir_list(input_gns, fail_time+1, token_fail)
        resp = response.json()
        if "code" in resp:
            if resp["code"] == 401001001:
                logger.error(resp)
                return self.get_dir_list(input_gns, fail_time+1, token_fail)
            else:
                logger.error(resp)
                return dict()

        return resp

    # 获取目录属性
    # @asd_decorator
    def get_dir_attribute(self, input_gns, fail_time=0, token_fail=False):
        if fail_time == 3:
            logger.error("get dir attribute fail exceed 3 times")
            return dict()
        print('开始获取AS token', __file__, 'get_dir_attribute')
        ret_code, token_id = asToken.get_token(self.ds_auth)
        url = '{}:{}/api/efast/v1/dir/attribute'.format(self.ds_address, self.ds_port)
        headers = {'Content-Type': 'application/json',
                   'Authorization': 'Bearer {}'.format(token_id)}
        payload = {
            "docid": input_gns
        }
        response = requests.request("POST", url, headers=headers, json=payload, timeout=30, verify=False)
        resp = response.json()
        if "code" in resp:
            if resp["code"] == 401001001:
                logger.error(resp)
                return self.get_dir_attribute(input_gns, fail_time+1, token_fail)
            else:
                logger.error(resp)
                return dict()

        return resp

    # 将gns 转换为path
    # @asd_decorator
    def get_convertpath(self, input_gns, fail_time=0, token_fail=False):
        if fail_time == 3:
            logger.error("get convertpath fail exceed 3 times")
            return dict()
        print('开始获取AS token', __file__, 'get_convertpath')
        ret_code, token_id = asToken.get_token(self.ds_auth)
        url = '{}:{}/api/efast/v1/file/convertpath'.format(self.ds_address, self.ds_port)
        headers = {'Content-Type': 'application/json',
                   'Authorization': 'Bearer {}'.format(token_id)}
        payload = {
            "docid": input_gns
        }
        response = requests.request("POST", url, headers=headers, json=payload, timeout=30, verify=False)
        resp = response.json()

        if "code" in resp:
            if resp["code"] == 401001001:
                print('开始获取AS token', __file__, 'get_convertpath 2')
                asToken.get_token(self.ds_auth)
                logger.error(resp)
                return self.get_convertpath(input_gns, fail_time+1, token_fail)
            else:
                logger.error(resp)
                return dict()

        return resp

    # 获取文档库信息
    # @asd_decorator
    def get_root_info(self, fail_time=0, token_fail=False):
        if fail_time == 3:
            logger.error("get root info fail exceed 3 times")
            return dict()
        print('开始获取AS token', __file__, 'get_root_info')
        ret_code, token_id = asToken.get_token(self.ds_auth)
        url = '{}:{}/api/efast/v1/entry-doc-lib'.format(self.ds_address, self.ds_port)
        headers = {'Content-Type': 'application/json',
                   'Authorization': 'Bearer {}'.format(token_id)}

        response = requests.request("GET", url, headers=headers, timeout=30, verify=False)
        resp = response.json()
        if "code" in resp:
            if resp["code"] == 401001001:
                logger.error(resp)
                return self.get_root_info(fail_time+1, token_fail)
            else:
                logger.error(resp)
                return dict()
        return resp

    def get_dir_info(self, dir_gns):
        path = ""
        dir_gns_split = dir_gns.split("/")
        doc_info_dict = dict()

        if len(dir_gns_split) > 3:
            last_dir_gns = "/".join(dir_gns_split[:-1])
            dir_info = self.get_dir_list(last_dir_gns)
            if len(dir_info) == 0:
                return {"error": "╮(╯▽╰)╭"}
            dir_info_dict = {r_info["docid"]: r_info for r_info in dir_info["dirs"]}
            doc_info_dict = {r_info["docid"]: r_info for r_info in dir_info["files"]}
            path = self.get_convertpath(last_dir_gns)["namepath"]
        else:
            dir_info = self.get_root_info()
            if len(dir_info) == 0:
                return {"error": "╮(╯▽╰)╭"}
            dir_info_dict = {r_info["id"]: r_info for r_info in dir_info}

        if dir_gns not in dir_info_dict:
            if dir_gns not in doc_info_dict:
                return {"error": "╮(╯▽╰)╭", "code": 0}
            doc_info_res = doc_info_dict[dir_gns]
            if "docid" in doc_info_res:
                doc_info_res["gns"] = doc_info_res["docid"]
            else:
                doc_info_res["creator"] = doc_info_res["created_by"]["name"]
                doc_info_res["editor"] = ""
                doc_info_res["gns"] = doc_info_res["id"]
            doc_info_res["path"] = path
            doc_info_res["root"] = 1
            doc_info_res["type"] = 0
            return {"error": "╮(╯▽╰)╭", "code": 1, "doc_info_res": doc_info_res}
        dir_info_res = dir_info_dict[dir_gns]
        if "docid" in dir_info_res:
            dir_info_res["gns"] = dir_info_res["docid"]
        else:
            dir_info_res["creator"] = dir_info_res["created_by"]["name"]
            dir_info_res["editor"] = ""
            dir_info_res["gns"] = dir_info_res["id"]
        dir_info_res["path"] = path
        dir_info_res["root"] = 1
        dir_info_res["type"] = 1

        return dir_info_res

    def compare_v1(self, input_info):
        gns = input_info["gns"]
        dir_info = self.get_dir_list(gns)
        child_folder_set = set(input_info["child_folder"])
        child_file_set = set(input_info["child_file"])

        dirs = dir_info["dirs"]
        dirs_dict = {dir["docid"]: dir for dir in dirs}
        files = dir_info["files"]
        files_dict = {file["docid"]: file for file in files}

        common_dir_gns = []
        add_dir_gns = []
        delete_dir_gns = []
        update_dir_gns = []
        rename_dir_gns = []
        for sg in child_folder_set:
            # 目录删除事件搜集
            if sg not in dirs_dict:
                self.change_collection.insert_one({
                    "gns": sg,
                    "type": DIR_TYPE,
                    "sub_root": 1,
                    "change_type": DELETE_STATUS
                })
                self.delete_path_collect(sg)
                delete_dir_gns.append(sg)
            else:
                common_dir_gns.append(sg)

        for sg, sg_info in dirs_dict.items():
            # 目录新增事件搜集
            if sg not in child_folder_set:

                sg_info["gns"] = sg_info["docid"]
                sg_info["type"] = DIR_TYPE
                sg_info["sub_root"] = 1
                sg_info["root_gns"] = input_info["root_gns"]
                sg_info["change_type"] = ADD_STATUS
                self.add_path_collect_v1(sg_info, input_info["path"])
                add_dir_gns.append(sg)

        for c_d_gns in common_dir_gns:
            ad_d_info = self.gns_info.find_one({"gns": c_d_gns})
            as_d_info = dirs_dict[c_d_gns]
            if ad_d_info["rev"] != as_d_info["rev"]:
                ad_d_info["rev"] = as_d_info["rev"]
                ad_d_info["type"] = DIR_TYPE
                ad_d_info["sub_root"] = 1
                ad_d_info["change_type"] = UPDATE_STATUS
                update_dir_gns.append(c_d_gns)
                self.compare_v1(ad_d_info)
            # 重命名事件搜集
            elif ad_d_info["name"] != as_d_info["name"]:
                ad_d_info["type"] = DIR_TYPE
                ad_d_info["sub_root"] = 1
                ad_d_info["change_type"] = RENAME_STATUS
                ad_d_info["name"] = as_d_info["name"]
                rename_path_list = ad_d_info["path"].split("/")
                rename_path_list[-1] = ad_d_info["name"]
                ad_d_info["path"] = "/".join(rename_path_list)
                update_insert_dir_info = {key: ad_d_info[key] for key in self.valid_change_key if key in ad_d_info}

                rename_dir_gns.append(c_d_gns)
                self.gns_info.update_many({"gns": c_d_gns}, {"$set": {"name": as_d_info["name"],
                                                                      "path": ad_d_info["path"]}})
                self.change_collection.insert_one(update_insert_dir_info)

        common_file_gns = []
        add_file_gns = []
        delete_file_gns = []
        update_file_gns = []
        rename_file_gns = []
        for sg in child_file_set:
            # 删除文件事件搜集
            if sg not in files_dict:
                self.change_collection.insert_one({
                    "gns": sg,
                    "type": FILE_TYPE,
                    "change_type": DELETE_STATUS
                })
                delete_file_gns.append(sg)
                # 暂时不动
            else:
                common_file_gns.append(sg)

        for sg, sg_info in files_dict.items():
            # 新增文件事件搜集
            if sg not in child_file_set:

                sg_info["gns"] = sg_info["docid"]
                sg_info["path"] = input_info["path"] + "/" + sg_info["name"]
                sg_info["change_type"] = ADD_STATUS
                sg_info["type"] = FILE_TYPE
                sg_info["root_gns"] = input_info["root_gns"]
                insert_file_info = {key: sg_info[key] for key in self.valid_key if key in sg_info}
                insert_insert_file_info = {key: sg_info[key] for key in self.valid_change_key if key in sg_info}
                add_file_gns.append(sg)
                self.gns_info.insert_one(insert_file_info)
                self.change_collection.insert_one(insert_insert_file_info)

        for c_f_gns in common_file_gns:
            ad_f_info = self.gns_info.find_one({"gns": c_f_gns})
            as_f_info = files_dict[c_f_gns]
            # 文件内容发生变化
            if ad_f_info["rev"] != as_f_info["rev"]:
                ad_f_info["rev"] = as_f_info["rev"]
                ad_f_info["name"] = as_f_info["name"]
                ad_f_info["change_type"] = UPDATE_STATUS
                update_file_gns.append(c_f_gns)
                update_insert_file_info = {key: ad_f_info[key] for key in self.valid_change_key if key in ad_f_info}
                self.gns_info.update_many({"gns": c_f_gns}, {"$set": {"name": ad_f_info["name"], "rev": ad_f_info["rev"]}})
                self.change_collection.insert_one(update_insert_file_info)
            # 文件名发生变化
            elif ad_f_info["name"] != as_f_info["name"]:

                ad_f_info["name"] = as_f_info["name"]
                ad_f_info["change_type"] = RENAME_STATUS
                rename_path_list = ad_f_info["path"].split("/")
                rename_path_list[-1] = ad_f_info["name"]
                ad_f_info["path"] = "/".join(rename_path_list)
                rename_file_gns.append(c_f_gns)
                update_insert_file_info = {key: ad_f_info[key] for key in self.valid_change_key if key in ad_f_info}
                self.gns_info.update_many({"gns": c_f_gns}, {"$set": {"name":  ad_f_info["name"],
                                                                      "path": ad_f_info["path"]}})
                self.change_collection.insert_one(update_insert_file_info)

        if add_dir_gns or delete_dir_gns or update_dir_gns or rename_dir_gns:
            input_info["child_folder"] = add_dir_gns + update_dir_gns + rename_dir_gns
        if add_file_gns or delete_file_gns or update_file_gns or rename_file_gns:
            input_info["child_file"] = add_file_gns + update_file_gns + rename_file_gns

        update_insert_dir_info = {key: input_info[key] for key in self.valid_change_key if key in input_info}
        self.change_collection.insert_one(update_insert_dir_info)
        self.gns_info.update_many({"gns": input_info["gns"]},
                                  {"$set": {"name":  input_info["name"],
                                           "rev": input_info["rev"],
                                           "child_folder": add_dir_gns + common_dir_gns,
                                           "child_file": add_file_gns + common_file_gns}})

    # 追踪所有因为删除目录所引起的删除事件
    def delete_path_collect(self, in_dir_gns):
        in_ad_dir_info = self.gns_info.find_one({"gns": in_dir_gns})
        for ifile_gns in in_ad_dir_info["child_file"]:
            self.change_collection.insert_one({
                "gns": ifile_gns,
                "type": FILE_TYPE,
                "change_type": DELETE_STATUS
            })

        for idir_gns in in_ad_dir_info["child_folder"]:
            self.change_collection.insert_one({
                "gns": idir_gns,
                "type": DIR_TYPE,
                "change_type": DELETE_STATUS
            })
            self.delete_path_collect(idir_gns)

    # 多线程模式
    def add_path_collect_v1(self, in_dir_info, path=""):
        self.queue.put((in_dir_info, path))
        while True:
            with ThreadPoolExecutor(max_workers=100) as executor:
                f_list = []
                while not self.queue.empty():
                    in_dir_info, path = self.queue.get()
                    th = executor.submit(self.add_path_collect_v2, in_dir_info, path)
                    f_list.append(th)

                for future in as_completed(f_list):
                    data = future.result()
                    logger.debug(data)
                    logger.debug('*' * 50)
                dir_count = 0
                logger.debug("thread num {}".format(len(f_list)))
                logger.debug(self.bk_queue.empty())
                while not self.bk_queue.empty():
                    dir_count += 1
                    self.queue.put(self.bk_queue.get())
                logger.debug("dir count {}".format(dir_count))
                if self.queue.empty():
                    break

    # 追踪所有因为添加目录所引起的添加事件
    def add_path_collect(self, in_dir_info, path=""):
        gns_info_list = []
        change_collection_list = []

        in_dir_gns = in_dir_info["gns"]
        in_dir_info["child_folder"] = []
        in_dir_info["child_file"] = []
        in_dir_info["child_folder_info"] = []
        in_dir_info["child_file_info"] = []

        in_dir_info["path"] = path + "/" + in_dir_info["name"] if path else in_dir_info["name"]
        in_as_dir_info = self.get_dir_list(in_dir_gns)
        for ifile in in_as_dir_info["files"]:
            ifile["gns"] = ifile["docid"]
            ifile["path"] = in_dir_info["path"] + "/" + ifile["name"]
            ifile["type"] = FILE_TYPE
            ifile["root_gns"] = in_dir_info["root_gns"]
            ifile["change_type"] = ADD_STATUS

            insert_file_info = {key: ifile[key] for key in self.valid_key if key in ifile}
            gns_info_list.append(insert_file_info)
            insert_change_file_info = {key: ifile[key] for key in self.valid_change_key if key in ifile}
            in_dir_info["child_file"].append(ifile["gns"])
            in_dir_info["child_file_info"].append({"name": ifile["name"], "gns": ifile["gns"]})

            change_collection_list.append(insert_change_file_info)

        for idir in in_as_dir_info["dirs"]:
            idir["gns"] = idir["docid"]
            idir["path"] = in_dir_info["path"] + "/" + idir["name"]
            idir["type"] = DIR_TYPE
            idir["change_type"] = ADD_STATUS
            idir["root_gns"] = in_dir_info["root_gns"]
            in_dir_info["child_folder"].append(idir["gns"])
            in_dir_info["child_folder_info"].append({"name": idir["name"], "gns": idir["gns"]})

            self.add_path_collect(idir, in_dir_info["path"])
        insert_dir_info = {key: in_dir_info[key] for key in self.valid_key if key in in_dir_info}
        gns_info_list.append(insert_dir_info)
        insert_change_dir_info = {key: in_dir_info[key] for key in self.valid_change_key if key in in_dir_info}
        change_collection_list.append(insert_change_dir_info)

        if len(gns_info_list):
            self.gns_info.insert_many(gns_info_list)
        if len(change_collection_list):
            self.change_collection.insert_many(change_collection_list)

    def insert_gnsinfo(self, info_list):
        self.gns_info.insert_many(info_list)

    def insert_changeinfo(self, info_list):
        self.change_collection.insert_many(info_list)

    # 多线程
    def add_path_collect_v2(self, in_dir_info, path=""):
        gns_info_list = []
        change_collection_list = []

        in_dir_gns = in_dir_info["gns"]
        in_dir_info["child_folder"] = []
        in_dir_info["child_file"] = []
        in_dir_info["child_folder_info"] = []
        in_dir_info["child_file_info"] = []

        in_dir_info["path"] = path + "/" + in_dir_info["name"] if path else in_dir_info["name"]
        logger.debug(in_dir_gns+" -> start")
        in_as_dir_info = self.get_dir_list(in_dir_gns)
        logger.debug(in_dir_gns+" -> end")
        logger.debug(in_as_dir_info)
        if len(in_as_dir_info) == 0:
            return "as has something wrong!"
        for ifile in in_as_dir_info["files"]:
            ifile["gns"] = ifile["docid"]
            ifile["path"] = in_dir_info["path"] + "/" + ifile["name"]
            ifile["type"] = FILE_TYPE
            ifile["root_gns"] = in_dir_info["root_gns"]
            ifile["change_type"] = ADD_STATUS

            insert_file_info = {key: ifile[key] for key in self.valid_key if key in ifile}
            gns_info_list.append(insert_file_info)
            insert_change_file_info = {key: ifile[key] for key in self.valid_change_key if key in ifile}
            in_dir_info["child_file"].append(ifile["gns"])
            in_dir_info["child_file_info"].append({"name": ifile["name"], "gns": ifile["gns"]})

            change_collection_list.append(insert_change_file_info)

        for idir in in_as_dir_info["dirs"]:
            idir["gns"] = idir["docid"]
            idir["path"] = in_dir_info["path"] + "/" + idir["name"]
            idir["type"] = DIR_TYPE
            idir["change_type"] = ADD_STATUS
            idir["root_gns"] = in_dir_info["root_gns"]
            in_dir_info["child_folder"].append(idir["gns"])
            in_dir_info["child_folder_info"].append({"name": idir["name"], "gns": idir["gns"]})

            self.bk_queue.put((idir, in_dir_info["path"]))

        insert_dir_info = {key: in_dir_info[key] for key in self.valid_key if key in in_dir_info}
        gns_info_list.append(insert_dir_info)
        insert_change_dir_info = {key: in_dir_info[key] for key in self.valid_change_key if key in in_dir_info}
        change_collection_list.append(insert_change_dir_info)

        self.gns_info.insert_many(gns_info_list)
        self.change_collection.insert_many(change_collection_list)

        return 1

    # 比较 anyshare 最新文档 和 anydata 所用文档的不同
    @print_run_time
    def compare(self):
        logger.info("anyshare gns info construct info start")
        self.change_collection.drop()

        # 是否有遗留问题
        last_change = 0
        for info in self.mongo_kg_db["{}_empty_data".format(self.graph_collection_prefix)].find({"root_gns": self.as7json["docid"]}, {"_id": 0}):
            last_change = 1
            self.change_collection.insert_one(info)
        self.mongo_kg_db["{}_empty_data".format(self.graph_collection_prefix)].remove({"root_gns": self.as7json["docid"]})
        print('开始获取AS token', __file__, 'compare')
        ret_token, obj_token = asToken.get_token(self.ds_auth)
        if ret_token != 200:
            logger.error(obj_token)
            raise Exception("anyshare token has something wrong! gns info extract fail!")
        ad_root_info = self.gns_info.find_one({"gns": self.as7json["docid"]})

        as_root_info = self.get_dir_info(self.as7json["docid"])
        if "error" in as_root_info:
            # 单个文件处理
            if as_root_info.get("code", 0) == 1:
                doc_info_res = as_root_info["doc_info_res"]
                update_insert_file_info = {key: doc_info_res[key] for key in self.valid_change_key if
                                           key in doc_info_res}
                if ad_root_info is None or len(ad_root_info) == 0:
                    self.gns_info.insert_one(update_insert_file_info)
                    update_insert_file_info["change_type"] = ADD_STATUS
                    self.change_collection.insert_one(update_insert_file_info)
                    return SOMETHING_CHANGE
                elif ad_root_info["rev"] == doc_info_res["rev"]:
                    if doc_info_res["name"] != ad_root_info["name"]:

                        self.gns_info.update_one({"gns": doc_info_res["gns"]}, {"$set": {"name":  doc_info_res["name"],
                                                                      "path": doc_info_res["path"]}})

                        self.change_collection.insert_one(update_insert_file_info)
                        return SOMETHING_CHANGE
                    if last_change:
                        return LASTESDATA_EMPTY
                    return NO_CHANGE
                else:
                    self.gns_info.update_many({"gns": doc_info_res["gns"]},
                                              {"$set": {"name": doc_info_res["name"], "rev": doc_info_res["rev"]}})
                    update_insert_file_info["change_type"] = UPDATE_STATUS
                    self.change_collection.insert_one(update_insert_file_info)
                    return SOMETHING_CHANGE

            logger.error(as_root_info)
            raise Exception("access anyshare root info fail!")

        # 这表明该数据源在as中被删
        if len(as_root_info) == 0:
            if ad_root_info is None or len(ad_root_info) == 0:
                return NO_CHANGE
            self.change_collection.insert_one({
                    "gns": ad_root_info["gns"],
                    "type": DIR_TYPE,
                    "root": 1,
                    "change_type": DELETE_STATUS
            })
            self.delete_path_collect(ad_root_info["gns"])
            return SOMETHING_CHANGE

        # 这表明该数据源在asz中新增
        if ad_root_info is None or len(ad_root_info) == 0:
            as_root_info["type"] = DIR_TYPE
            as_root_info["change_type"] = ADD_STATUS
            as_root_info["root_gns"] = as_root_info["gns"]
            self.add_path_collect_v1(as_root_info, as_root_info["path"])
            return SOMETHING_CHANGE

        if ad_root_info["rev"] == as_root_info["rev"]:
            # 如果名字也是一样的，那么就没有变化
            if ad_root_info["name"] == as_root_info["name"]:
                if last_change:
                    return LASTESDATA_EMPTY
                return NO_CHANGE
            else:
                ad_root_info["type"] = DIR_TYPE
                ad_root_info["root"] = 1
                ad_root_info["change_type"] = RENAME_STATUS
                ad_root_info["name"] = as_root_info["name"]
                rename_path_list = ad_root_info["path"].split("/")
                rename_path_list[-1] = ad_root_info["name"]
                ad_root_info["path"] = "/".join(rename_path_list)
                update_insert_dir_info = {key: ad_root_info[key] for key in self.valid_change_key if key in ad_root_info}

                self.gns_info.update_many({"gns": ad_root_info}, {"$set": {"name": ad_root_info["name"],
                                                                      "path": ad_root_info["path"]}})
                self.change_collection.insert_one(update_insert_dir_info)

                return SOMETHING_CHANGE

        ad_root_info["rev"] = as_root_info["rev"]
        ad_root_info["change_type"] = UPDATE_STATUS
        ad_root_info["type"] = DIR_TYPE
        # 深入比较
        self.compare_v1(ad_root_info)

        return SOMETHING_CHANGE

    #  构建文档树
    def build_tree(self, root_gns, path=""):
        self.gns_info.drop()
        all_root_info = self.get_root_info()
        all_root_dict = {r_info["id"]: r_info for r_info in all_root_info}
        root_info = all_root_dict[root_gns]
        root_info["gns"] = root_info.pop("id")
        root_info["root"] = 1
        root_info["type"] = 1
        queue_list = []
        current_path = root_info["name"]
        if path:
            current_path = path + "/" + root_info["name"]
        dir_info = self.get_dir_list(root_gns)
        for dir in dir_info["dirs"]:
            queue_list.append((dir, current_path))

        root_info["path"] = current_path
        root_info["child_folder"] = [ele["docid"] for ele in dir_info["dirs"]]
        root_info["child_file"] = [ele["docid"] for ele in dir_info["files"]]
        root_info["creator"] = root_info["created_by"]["name"]
        root_info["editor"] = ""

        self.gns_info.insert_one(root_info)
        for file in dir_info["files"]:
            file["gns"] = file["docid"]
            file["type"] = 0
            file["path"] = current_path + "/" + file["name"]
        if dir_info["files"]:
            self.gns_info.insert_many(dir_info["files"])

        # 广度优先的方式遍历
        while queue_list:
            q_info, last_path = queue_list.pop(0)
            q_info["gns"] = q_info["docid"]
            q_info["path"] = last_path + "/" + q_info["name"]
            q_info["type"] = 1

            sub_dir_info = self.get_dir_list(q_info["gns"])
            for dir in sub_dir_info["dirs"]:
                queue_list.append((dir, q_info["path"]))

            q_info["child_folder"] = [ele["docid"] for ele in sub_dir_info["dirs"]]
            q_info["child_file"] = [ele["docid"] for ele in sub_dir_info["files"]]
            self.gns_info.insert_one(q_info)

            for sub_file in sub_dir_info["files"]:
                sub_file["gns"] = sub_file["docid"]
                sub_file["path"] = q_info["path"] + "/" + sub_file["name"]
                sub_file["type"] = 0
            if sub_dir_info["files"]:
                self.gns_info.insert_many(sub_dir_info["files"])

    def add_change_info(self):
        change_list = self.change_collection.find({"change_type": {"$in": [UPDATE_STATUS, ADD_STATUS, RENAME_STATUS]}})
        for item in change_list:
            gns_info = self.gns_info.find_one({"gns": item["gns"]}, {"_id": 0, "gns": 0})
            self.change_collection.update_many({"gns": item["gns"]}, {"$set": gns_info})
